import os
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
businesses = [
    {
        "id": "BIZ-01",
        "name": "Leopold Cafe & Bar",
        "category": "Restaurants & Cafes",
        "location": "Mumbai",
        "rating": 4.2,
        "reviews": 8400,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "http://leopoldcafe.com",
        "https_redirect": "No",
        "ssl_status": "None",
        "ssl_provider": "None",
        "admin_panel_exposed": "Yes",
        "forms_protection": "Unprotected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "X-Frame-Options, CSP, HSTS",
    },
    {
        "id": "BIZ-02",
        "name": "Toit Brewpub",
        "category": "Restaurants & Cafes",
        "location": "Bengaluru",
        "rating": 4.6,
        "reviews": 9800,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://toit.in",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "Let's Encrypt",
        "admin_panel_exposed": "Yes",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "CSP, HSTS",
    },
    {
        "id": "BIZ-03",
        "name": "Glen's Bakehouse",
        "category": "Restaurants & Cafes",
        "location": "Bengaluru",
        "rating": 4.4,
        "reviews": 3200,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://glensbakehouse.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "Cloudflare",
        "admin_panel_exposed": "No",
        "forms_protection": "Unprotected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-04",
        "name": "Flurys",
        "category": "Restaurants & Cafes",
        "location": "Kolkata",
        "rating": 4.3,
        "reviews": 4500,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "http://flurys.com",
        "https_redirect": "No",
        "ssl_status": "None",
        "ssl_provider": "None",
        "admin_panel_exposed": "No",
        "forms_protection": "Unprotected",
        "directory_listing": "Yes",
        "safe_browsing": "Clean",
        "missing_headers": "All Headers Missing",
    },
    {
        "id": "BIZ-05",
        "name": "Mavalli Tiffin Room (MTR)",
        "category": "Restaurants & Cafes",
        "location": "Bengaluru",
        "rating": 4.5,
        "reviews": 12000,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://mtrrestaurants.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-06",
        "name": "Paradise Biryani",
        "category": "Restaurants & Cafes",
        "location": "Hyderabad",
        "rating": 4.0,
        "reviews": 18500,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://paradisebiryani.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "Sectigo",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "X-Frame-Options",
    },
    {
        "id": "BIZ-07",
        "name": "Peter Cat",
        "category": "Restaurants & Cafes",
        "location": "Kolkata",
        "rating": 4.5,
        "reviews": 8900,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "http://petercatkolkata.co.in",
        "https_redirect": "No",
        "ssl_status": "Expired",
        "ssl_provider": "Let's Encrypt (Expired)",
        "admin_panel_exposed": "Yes",
        "forms_protection": "Unprotected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "All Headers Missing",
    },
    {
        "id": "BIZ-08",
        "name": "Cafe Goodluck",
        "category": "Restaurants & Cafes",
        "location": "Pune",
        "rating": 4.3,
        "reviews": 6200,
        "gbp_claimed": "No",
        "website_status": "No Website",
        "website_url": "None",
        "https_redirect": "N/A",
        "ssl_status": "None",
        "ssl_provider": "None",
        "admin_panel_exposed": "N/A",
        "forms_protection": "N/A",
        "directory_listing": "N/A",
        "safe_browsing": "N/A",
        "missing_headers": "N/A",
    },
    {
        "id": "BIZ-09",
        "name": "Britto's Restaurant & Bar",
        "category": "Restaurants & Cafes",
        "location": "Goa",
        "rating": 4.1,
        "reviews": 11500,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "http://brittosgoa.com",
        "https_redirect": "No",
        "ssl_status": "None",
        "ssl_provider": "None",
        "admin_panel_exposed": "No",
        "forms_protection": "Unprotected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "X-Frame-Options, CSP",
    },
    {
        "id": "BIZ-10",
        "name": "Karim's Restaurant",
        "category": "Restaurants & Cafes",
        "location": "Delhi",
        "rating": 4.2,
        "reviews": 5600,
        "gbp_claimed": "No",
        "website_status": "Active",
        "website_url": "http://karimsdelhi.in",
        "https_redirect": "No",
        "ssl_status": "None",
        "ssl_provider": "None",
        "admin_panel_exposed": "Yes",
        "forms_protection": "Unprotected",
        "directory_listing": "Yes",
        "safe_browsing": "Clean",
        "missing_headers": "All Headers Missing",
    },
    {
        "id": "BIZ-11",
        "name": "Apollo Hospitals Greams Road",
        "category": "Healthcare & Clinics",
        "location": "Chennai",
        "rating": 4.2,
        "reviews": 4200,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://apollohospitals.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-12",
        "name": "Fortis Hospital Mulund",
        "category": "Healthcare & Clinics",
        "location": "Mumbai",
        "rating": 4.1,
        "reviews": 2800,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://fortishealthcare.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "Entrust",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-13",
        "name": "Manipal Hospital Old Airport Rd",
        "category": "Healthcare & Clinics",
        "location": "Bengaluru",
        "rating": 4.3,
        "reviews": 3500,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://manipalhospitals.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "Sectigo",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-14",
        "name": "Lilavati Hospital & Research Centre",
        "category": "Healthcare & Clinics",
        "location": "Mumbai",
        "rating": 4.0,
        "reviews": 1900,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://lilavatihospital.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "GlobalSign",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "CSP",
    },
    {
        "id": "BIZ-15",
        "name": "Narayana Health City",
        "category": "Healthcare & Clinics",
        "location": "Bengaluru",
        "rating": 4.4,
        "reviews": 3100,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://narayanahealth.org",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-16",
        "name": "Ganga Hospital",
        "category": "Healthcare & Clinics",
        "location": "Coimbatore",
        "rating": 4.6,
        "reviews": 1500,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "http://gangahospital.co.in",
        "https_redirect": "No",
        "ssl_status": "None",
        "ssl_provider": "None",
        "admin_panel_exposed": "Yes",
        "forms_protection": "Unprotected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "All Headers Missing",
    },
    {
        "id": "BIZ-17",
        "name": "Medanta The Medicity",
        "category": "Healthcare & Clinics",
        "location": "Delhi",
        "rating": 4.3,
        "reviews": 8900,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://medanta.org",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-18",
        "name": "KIMS Hospitals Secunderabad",
        "category": "Healthcare & Clinics",
        "location": "Hyderabad",
        "rating": 4.1,
        "reviews": 2400,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://kimshospitals.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "Sectigo",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-19",
        "name": "Dr. Mehta's Hospitals",
        "category": "Healthcare & Clinics",
        "location": "Chennai",
        "rating": 3.9,
        "reviews": 950,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "http://mehtahospitals.com",
        "https_redirect": "No",
        "ssl_status": "Expired",
        "ssl_provider": "Let's Encrypt (Expired)",
        "admin_panel_exposed": "No",
        "forms_protection": "Unprotected",
        "directory_listing": "No",
        "safe_browsing": "Suspicious",
        "missing_headers": "All Headers Missing",
    },
    {
        "id": "BIZ-20",
        "name": "City Dental Clinic & Implant Center",
        "category": "Healthcare & Clinics",
        "location": "Ahmedabad",
        "rating": 4.2,
        "reviews": 120,
        "gbp_claimed": "No",
        "website_status": "No Website",
        "website_url": "None",
        "https_redirect": "N/A",
        "ssl_status": "None",
        "ssl_provider": "None",
        "admin_panel_exposed": "N/A",
        "forms_protection": "N/A",
        "directory_listing": "N/A",
        "safe_browsing": "N/A",
        "missing_headers": "N/A",
    },
    {
        "id": "BIZ-21",
        "name": "Nalli Silks T-Nagar",
        "category": "Retail & Shops",
        "location": "Chennai",
        "rating": 4.5,
        "reviews": 3200,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://nalli.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "GeoTrust",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "CSP",
    },
    {
        "id": "BIZ-22",
        "name": "Bahrisons Booksellers Khan Mkt",
        "category": "Retail & Shops",
        "location": "Delhi",
        "rating": 4.7,
        "reviews": 890,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://bahrisons.in",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "Sectigo",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "CSP, HSTS",
    },
    {
        "id": "BIZ-23",
        "name": "Sabyasachi Flagship Store Kolkata",
        "category": "Retail & Shops",
        "location": "Kolkata",
        "rating": 4.6,
        "reviews": 650,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://sabyasachi.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-24",
        "name": "Sarojini Saree House",
        "category": "Retail & Shops",
        "location": "Jaipur",
        "rating": 3.6,
        "reviews": 45,
        "gbp_claimed": "No",
        "website_status": "No Website",
        "website_url": "None",
        "https_redirect": "N/A",
        "ssl_status": "None",
        "ssl_provider": "None",
        "admin_panel_exposed": "N/A",
        "forms_protection": "N/A",
        "directory_listing": "N/A",
        "safe_browsing": "N/A",
        "missing_headers": "N/A",
    },
    {
        "id": "BIZ-25",
        "name": "Mysore Silk Udyog",
        "category": "Retail & Shops",
        "location": "Bengaluru",
        "rating": 4.1,
        "reviews": 1200,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "http://mysoresilkudyog.com",
        "https_redirect": "No",
        "ssl_status": "None",
        "ssl_provider": "None",
        "admin_panel_exposed": "No",
        "forms_protection": "Unprotected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "X-Frame-Options",
    },
    {
        "id": "BIZ-26",
        "name": "Fabindia Fort",
        "category": "Retail & Shops",
        "location": "Mumbai",
        "rating": 4.2,
        "reviews": 980,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://fabindia.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-27",
        "name": "Manyavar",
        "category": "Retail & Shops",
        "location": "Kolkata",
        "rating": 4.4,
        "reviews": 1500,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://manyavar.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "Sectigo",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-28",
        "name": "Kalaniketan Colaba",
        "category": "Retail & Shops",
        "location": "Mumbai",
        "rating": 4.0,
        "reviews": 720,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "http://kalaniketancolaba.in",
        "https_redirect": "No",
        "ssl_status": "None",
        "ssl_provider": "None",
        "admin_panel_exposed": "Yes",
        "forms_protection": "Unprotected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "All Headers Missing",
    },
    {
        "id": "BIZ-29",
        "name": "Ratnadeep Supermarket Jubilee Hills",
        "category": "Retail & Shops",
        "location": "Hyderabad",
        "rating": 4.3,
        "reviews": 1100,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://ratnadeep.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "Cloudflare",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "CSP",
    },
    {
        "id": "BIZ-30",
        "name": "Coimbatore Grocery Mart",
        "category": "Retail & Shops",
        "location": "Coimbatore",
        "rating": 3.7,
        "reviews": 35,
        "gbp_claimed": "No",
        "website_status": "No Website",
        "website_url": "None",
        "https_redirect": "N/A",
        "ssl_status": "None",
        "ssl_provider": "None",
        "admin_panel_exposed": "N/A",
        "forms_protection": "N/A",
        "directory_listing": "N/A",
        "safe_browsing": "N/A",
        "missing_headers": "N/A",
    },
    {
        "id": "BIZ-31",
        "name": "Cyril Amarchand Mangaldas",
        "category": "Professional Services",
        "location": "Mumbai",
        "rating": 4.5,
        "reviews": 210,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://cyrilshroff.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "CSP, HSTS",
    },
    {
        "id": "BIZ-32",
        "name": "AZB & Partners Corporate Office",
        "category": "Professional Services",
        "location": "Mumbai",
        "rating": 4.6,
        "reviews": 180,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://azbpartners.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-33",
        "name": "Trilegal Advocates",
        "category": "Professional Services",
        "location": "Delhi",
        "rating": 4.4,
        "reviews": 150,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://trilegal.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "Sectigo",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "CSP",
    },
    {
        "id": "BIZ-34",
        "name": "Khaitan & Co Advocates",
        "category": "Professional Services",
        "location": "Kolkata",
        "rating": 4.5,
        "reviews": 190,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://khaitanco.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-35",
        "name": "Luthra & Luthra Law Offices Delhi",
        "category": "Professional Services",
        "location": "Delhi",
        "rating": 4.3,
        "reviews": 140,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://luthra.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-36",
        "name": "J. Sagar Associates",
        "category": "Professional Services",
        "location": "Delhi",
        "rating": 4.2,
        "reviews": 120,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://jsalaw.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "Sectigo",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-37",
        "name": "Singhania & Co. Legal Advisors",
        "category": "Professional Services",
        "location": "Bengaluru",
        "rating": 4.1,
        "reviews": 85,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://singhania.in",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "Sectigo",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "CSP",
    },
    {
        "id": "BIZ-38",
        "name": "Gupta Chartered Accountants",
        "category": "Professional Services",
        "location": "Pune",
        "rating": 3.8,
        "reviews": 12,
        "gbp_claimed": "No",
        "website_status": "No Website",
        "website_url": "None",
        "https_redirect": "N/A",
        "ssl_status": "None",
        "ssl_provider": "None",
        "admin_panel_exposed": "N/A",
        "forms_protection": "N/A",
        "directory_listing": "N/A",
        "safe_browsing": "N/A",
        "missing_headers": "N/A",
    },
    {
        "id": "BIZ-39",
        "name": "Mehta Tax Advisors",
        "category": "Professional Services",
        "location": "Mumbai",
        "rating": 3.7,
        "reviews": 18,
        "gbp_claimed": "No",
        "website_status": "No Website",
        "website_url": "None",
        "https_redirect": "N/A",
        "ssl_status": "None",
        "ssl_provider": "None",
        "admin_panel_exposed": "N/A",
        "forms_protection": "N/A",
        "directory_listing": "N/A",
        "safe_browsing": "N/A",
        "missing_headers": "N/A",
    },
    {
        "id": "BIZ-40",
        "name": "Sharma Real Estate Consultants",
        "category": "Professional Services",
        "location": "Jaipur",
        "rating": 3.5,
        "reviews": 8,
        "gbp_claimed": "No",
        "website_status": "No Website",
        "website_url": "None",
        "https_redirect": "N/A",
        "ssl_status": "None",
        "ssl_provider": "None",
        "admin_panel_exposed": "N/A",
        "forms_protection": "N/A",
        "directory_listing": "N/A",
        "safe_browsing": "N/A",
        "missing_headers": "N/A",
    },
    {
        "id": "BIZ-41",
        "name": "Taj Mahal Palace Colaba",
        "category": "Hospitality & Hotels",
        "location": "Mumbai",
        "rating": 4.9,
        "reviews": 18500,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://tajhotels.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-42",
        "name": "The Oberoi Mahatma Gandhi Rd",
        "category": "Hospitality & Hotels",
        "location": "Bengaluru",
        "rating": 4.8,
        "reviews": 6200,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://oberoihotels.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-43",
        "name": "The Leela Palace MRC Nagar",
        "category": "Hospitality & Hotels",
        "location": "Chennai",
        "rating": 4.7,
        "reviews": 4800,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://theleela.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-44",
        "name": "ITC Grand Chola Guindy",
        "category": "Hospitality & Hotels",
        "location": "Chennai",
        "rating": 4.8,
        "reviews": 9500,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://itchotels.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "Sectigo",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-45",
        "name": "JW Marriott Senapati Bapat Rd",
        "category": "Hospitality & Hotels",
        "location": "Pune",
        "rating": 4.7,
        "reviews": 5400,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://marriott.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-46",
        "name": "Rambagh Palace",
        "category": "Hospitality & Hotels",
        "location": "Jaipur",
        "rating": 4.9,
        "reviews": 3200,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://tajhotels.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-47",
        "name": "Wildflower Hall Anamika",
        "category": "Hospitality & Hotels",
        "location": "Shimla",
        "rating": 4.8,
        "reviews": 1100,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://oberoihotels.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "DigiCert",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-48",
        "name": "Brunton Boatyard",
        "category": "Hospitality & Hotels",
        "location": "Kochi",
        "rating": 4.5,
        "reviews": 450,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "https://cghearth.com",
        "https_redirect": "Yes",
        "ssl_status": "Valid",
        "ssl_provider": "GeoTrust",
        "admin_panel_exposed": "No",
        "forms_protection": "Protected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "None",
    },
    {
        "id": "BIZ-49",
        "name": "Gokulam Park Plaza",
        "category": "Hospitality & Hotels",
        "location": "Kochi",
        "rating": 4.0,
        "reviews": 750,
        "gbp_claimed": "Yes",
        "website_status": "Active",
        "website_url": "http://gokulampark.co.in",
        "https_redirect": "No",
        "ssl_status": "Expired",
        "ssl_provider": "Let's Encrypt (Expired)",
        "admin_panel_exposed": "No",
        "forms_protection": "Unprotected",
        "directory_listing": "No",
        "safe_browsing": "Clean",
        "missing_headers": "All Headers Missing",
    },
    {
        "id": "BIZ-50",
        "name": "Shimla Residency Guest House",
        "category": "Hospitality & Hotels",
        "location": "Shimla",
        "rating": 3.5,
        "reviews": 45,
        "gbp_claimed": "No",
        "website_status": "No Website",
        "website_url": "None",
        "https_redirect": "N/A",
        "ssl_status": "None",
        "ssl_provider": "None",
        "admin_panel_exposed": "N/A",
        "forms_protection": "N/A",
        "directory_listing": "N/A",
        "safe_browsing": "N/A",
        "missing_headers": "N/A",
    }
]
os.makedirs("charts", exist_ok=True)
processed_data = []
for biz in businesses:
    row = biz.copy()
    risk_indicators_list = []
    if row["website_status"] == "No Website":
        if row["gbp_claimed"] == "No":
            row["security_score"] = 0
            row["key_security_gaps"] = "Digital Absence, Unclaimed Google Business Profile"
        else:
            row["security_score"] = 1
            row["key_security_gaps"] = "Digital Absence, No Website Present"
        risk_indicators_list.append("No Website (High Risk)")
    else:
        score = 5
        gaps = []
        if row["https_redirect"] == "No":
            score -= 1
            gaps.append("No HTTPS Redirect")
            risk_indicators_list.append("No HTTPS / Not Secure")
        if row["ssl_status"] in ["Expired", "None"]:
            score -= 1
            gaps.append(f"SSL Status: {row['ssl_status']}")
            risk_indicators_list.append("SSL Certificate Expired/Invalid" if row["ssl_status"] == "Expired" else "No SSL Certificate")
        if row["admin_panel_exposed"] == "Yes":
            score -= 1
            gaps.append("Exposed Admin Panel (/wp-admin or /admin)")
            risk_indicators_list.append("Admin Panel Exposed")
        if row["forms_protection"] == "Unprotected" or row["directory_listing"] == "Yes":
            score -= 1
            subgaps = []
            if row["forms_protection"] == "Unprotected":
                subgaps.append("Unprotected Forms")
                risk_indicators_list.append("Forms Without Protection")
            if row["directory_listing"] == "Yes":
                subgaps.append("Directory Listing Enabled")
                risk_indicators_list.append("Open Directory Exposed")
            gaps.append(", ".join(subgaps))
        if row["safe_browsing"] in ["Suspicious", "Blacklisted"]:
            score -= 1
            gaps.append(f"Safe Browsing: {row['safe_browsing']}")
            risk_indicators_list.append("Malware/Phishing Risk Detected")
        if score < 1:
            score = 1
        row["security_score"] = score
        row["key_security_gaps"] = ", ".join(gaps) if gaps else "None (All Secure)"
    if row["security_score"] == 5:
        row["risk_level"] = "Very Secure"
    elif row["security_score"] == 4:
        row["risk_level"] = "Secure"
    elif row["security_score"] == 3:
        row["risk_level"] = "Moderate Risk"
    elif row["security_score"] == 2:
        row["risk_level"] = "High Risk"
    else:
        row["risk_level"] = "Very High Risk"
    row["risk_indicators"] = risk_indicators_list
    rec = ""
    if row["website_status"] == "No Website":
        if row["gbp_claimed"] == "No":
            rec = "1. Claim Google Business Profile immediately to prevent hijacking. 2. Build a secure, mobile-friendly landing page with HTTPS. 3. Standardize NAP (Name, Address, Phone) details."
        else:
            rec = "1. Invest in a basic secure website (HTTPS) to capture organic search traffic. 2. Provide customer inquiries via webforms instead of just phone. 3. Setup official email domain."
    else:
        recs = []
        if row["https_redirect"] == "No" or row["ssl_status"] in ["Expired", "None"]:
            recs.append("Install valid SSL certificate (e.g. Let's Encrypt) and force HTTPS redirect")
        if row["admin_panel_exposed"] == "Yes":
            recs.append("Restrict access to admin pages (/wp-admin, /admin) by IP address, change default URL path, and enable MFA")
        if row["directory_listing"] == "Yes":
            recs.append("Disable directory indexing in server config (.htaccess or Nginx config)")
        if row["forms_protection"] == "Unprotected":
            recs.append("Implement CSRF tokens, secure input sanitization, and captchas on contact forms")
        if row["safe_browsing"] in ["Suspicious", "Blacklisted"]:
            recs.append("Run security malware scans, remove malicious scripts, and request Google review of blacklist")
        if not recs:
            recs.append("Maintain existing secure posture. Schedule routine security auditing and patch CMS regularly.")
        rec = "; ".join(recs)
    row["recommendation"] = rec
    if row["security_score"] <= 2:
        row["outreach_priority"] = "CRITICAL"
    elif row["security_score"] == 3:
        row["outreach_priority"] = "MEDIUM"
    else:
        row["outreach_priority"] = "LOW"
    if row["website_status"] == "No Website":
        if row["gbp_claimed"] == "No":
            row["sales_pitch_hook"] = "Profile Hijacking Risk (Unclaimed Google Maps Profile) & Digital Absence"
        else:
            row["sales_pitch_hook"] = "Digital Absence (Missing Website Customer Acquisition Channel)"
    else:
        gaps_pitch = []
        if row["ssl_status"] in ["Expired", "None"] or row["https_redirect"] == "No":
            gaps_pitch.append("Insecure Protocol / Missing SSL (Active Browser Security Warnings)")
        if row["admin_panel_exposed"] == "Yes":
            gaps_pitch.append("Exposed WordPress Admin Panel (Brute Force Vulnerability)")
        if row["gbp_claimed"] == "No":
            gaps_pitch.append("Unclaimed Business Profile on Google Maps (Malicious Takeover Risk)")
        if row["directory_listing"] == "Yes":
            gaps_pitch.append("Exposed Web Directory (Insecure Server File Listings)")
        if row["forms_protection"] == "Unprotected":
            gaps_pitch.append("Insecure Webforms (Vulnerable to SQLi and Contact Spam)")
        if row["safe_browsing"] in ["Suspicious", "Blacklisted"]:
            gaps_pitch.append("Phishing/Malware Flags Detected on Domain")
        if gaps_pitch:
            row["sales_pitch_hook"] = " | ".join(gaps_pitch)
        else:
            row["sales_pitch_hook"] = "General Maintenance & Continuous Security Penetration Scanning"
    if row["security_score"] <= 2:
        if row["reviews"] >= 1000:
            row["estimated_lead_value"] = "$3,500 - $5,000 (Premium Retainer)"
        elif row["reviews"] >= 100:
            row["estimated_lead_value"] = "$1,500 - $2,500 (Standard Setup)"
        else:
            row["estimated_lead_value"] = "$750 - $1,200 (Essential Setup)"
    elif row["security_score"] == 3:
        if row["reviews"] >= 100:
            row["estimated_lead_value"] = "$1,000 - $1,800 (Remediation Package)"
        else:
            row["estimated_lead_value"] = "$500 - $1,000 (Patch Package)"
    else:
        row["estimated_lead_value"] = "$250 (Routine Audit / Audit Retainer)"
    if row["website_status"] == "No Website":
        row["outreach_strategy"] = "Local Business Visit (Claim Google Maps Listing Pitch)"
    elif row["reviews"] >= 1000 and row["security_score"] <= 2:
        row["outreach_strategy"] = "In-Person Executive Meeting (Deliver Printed Security Briefing)"
    elif row["security_score"] <= 2:
        row["outreach_strategy"] = "Direct Email Outreach + Safe Browsing Alert PDF"
    else:
        row["outreach_strategy"] = "LinkedIn Outreach to Owner / Operations Manager"
    name_clean = row["name"]
    if row["ssl_status"] == "Expired":
        row["outreach_subject_line"] = f"Urgent: Expired SSL Certificate security warning for {name_clean}"
    elif row["ssl_status"] == "None" or row["https_redirect"] == "No":
        if row["website_status"] != "No Website":
            row["outreach_subject_line"] = f"Critical Security Warning: Insecure connection on {name_clean}"
        else:
            row["outreach_subject_line"] = f"Unsecured local presence for {name_clean}"
    elif row["admin_panel_exposed"] == "Yes":
        row["outreach_subject_line"] = f"Security Risk Alert: Exposed WordPress admin panel on {name_clean}"
    elif row["gbp_claimed"] == "No":
        row["outreach_subject_line"] = f"Action Required: Unclaimed Google Maps Profile for {name_clean}"
    else:
        row["outreach_subject_line"] = f"Digital Security Posture Review for {name_clean}"
    if row["website_status"] == "No Website":
        row["first_touch_pitch"] = f"Your Google Maps profile for {name_clean} is currently unclaimed, meaning anyone can change your phone number or edit your address. We can secure it today."
    elif row["ssl_status"] == "Expired":
        row["first_touch_pitch"] = f"Your website ({row['website_url']}) currently triggers a red 'Your Connection is Not Private' warning in browsers. We can restore a secure SSL certificate within 2 hours."
    elif row["ssl_status"] == "None" or row["https_redirect"] == "No":
        row["first_touch_pitch"] = f"Your website ({row['website_url']}) is sending customer passwords and inputs in cleartext without forcing HTTPS, prompting security warnings. We can fix it today."
    elif row["admin_panel_exposed"] == "Yes":
        row["first_touch_pitch"] = f"Your administrative login page is accessible to anyone at {row['website_url']}/wp-admin, leaving you open to brute-force credential stuffing. We can restrict this immediately."
    elif row["gbp_claimed"] == "No":
        row["first_touch_pitch"] = f"While your website is secure, your Google Maps listing is currently unclaimed, exposing you to profile hijacking. We can lock down your Maps listing."
    else:
        row["first_touch_pitch"] = f"We have run a digital presence scan on {name_clean} and verified that your primary protocols are secure. We offer routine scanning packages to maintain this posture."
    processed_data.append(row)
processed_data.sort(key=lambda x: (x["security_score"], -x["reviews"]))
df = pd.DataFrame(processed_data)
csv_columns_order = [
    "id", "name", "category", "location", "rating", "reviews", "gbp_claimed", 
    "website_status", "website_url", "https_redirect", "ssl_status", "ssl_provider", 
    "admin_panel_exposed", "forms_protection", "directory_listing", "safe_browsing", 
    "missing_headers", "security_score", "risk_level", "outreach_priority", 
    "estimated_lead_value", "outreach_strategy", "outreach_subject_line", 
    "first_touch_pitch", "sales_pitch_hook", "key_security_gaps", "recommendation"
]
import csv
df_csv = df[csv_columns_order].copy()
df_csv.columns = [
    col.replace("_", " ").title()
    .replace("Id", "ID")
    .replace("Gbp", "GBP")
    .replace("Https", "HTTPS")
    .replace("Ssl", "SSL")
    for col in df_csv.columns
]
df_csv.to_csv("National_Digital_Risk_Audit_Database.csv", index=False, quoting=csv.QUOTE_ALL)
print("Processed security data saved to National_Digital_Risk_Audit_Database.csv")
print("Configuring dark theme for charts...")
def setup_dark_matplotlib_theme():
    plt.rcParams['figure.facecolor'] = '#120F1D'
    plt.rcParams['axes.facecolor'] = '#120F1D'
    plt.rcParams['text.color'] = '#FAF6F0'
    plt.rcParams['axes.labelcolor'] = '#FAF6F0'
    plt.rcParams['xtick.color'] = '#D5C3C6'
    plt.rcParams['ytick.color'] = '#D5C3C6'
    plt.rcParams['grid.color'] = '#312C45' 
    plt.rcParams['grid.alpha'] = 0.5
    plt.rcParams['font.family'] = 'sans-serif'
    plt.rcParams['font.sans-serif'] = ['Segoe UI', 'Arial', 'DejaVu Sans']
    plt.rcParams['font.size'] = 12
    plt.rcParams['axes.titlesize'] = 16
    plt.rcParams['axes.labelsize'] = 14
    plt.rcParams['xtick.labelsize'] = 12
    plt.rcParams['ytick.labelsize'] = 12
    plt.rcParams['legend.fontsize'] = 11
setup_dark_matplotlib_theme()
cyber_colors = {
    "Very Secure": "#88D49E",    
    "Secure": "#2E8B57",         
    "Moderate Risk": "#F1C40F",   
    "High Risk": "#E67E22",      
    "Very High Risk": "#F07167"  
}
print("Generating distribution pie chart...")
risk_counts = df["risk_level"].value_counts().reindex(["Very Secure", "Secure", "Moderate Risk", "High Risk", "Very High Risk"]).fillna(0)
risk_counts = risk_counts[risk_counts > 0] 
fig, ax = plt.subplots(figsize=(7, 6))
colors = [cyber_colors[label] for label in risk_counts.index]
wedges, texts, autotexts = ax.pie(
    risk_counts, 
    autopct='%1.0f%%', 
    startangle=140, 
    colors=colors,
    wedgeprops=dict(width=0.45, edgecolor='#120F1D', linewidth=4),
    pctdistance=0.75  
)
for autotext in autotexts:
    autotext.set_color('#111827')
    autotext.set_weight('bold')
    autotext.set_fontsize(12)
ax.text(0, 0, f'{len(df)}\nTargets', ha='center', va='center', fontsize=20, fontweight='bold', color='#E5A88B')
ax.legend(wedges, risk_counts.index, title="Risk Level", loc="center left", bbox_to_anchor=(0.95, 0.5), 
          facecolor='#1B1829', edgecolor='#312C45', labelcolor='#FAF6F0')
plt.title(f"Overall Security Risk Distribution\n({len(df)} Businesses Checked)", fontsize=16, weight="bold", pad=20, color="#E5A88B")
plt.savefig("charts/risk_distribution.png", dpi=300, bbox_inches='tight', facecolor='#120F1D')
plt.close()
print("Generating category scores bar chart...")
cat_scores = df.groupby("category")["security_score"].mean().sort_values(ascending=True)
fig, ax = plt.subplots(figsize=(8, 4.5))
palette = ["#F07167", "#E67E22", "#F1C40F", "#D5C3C6", "#E5A88B"] 
bars = ax.barh(cat_scores.index, cat_scores.values, color=palette, height=0.55, edgecolor='#312C45', alpha=0.95)
plt.xlabel("Average Security Score (0 - 5)", fontsize=12, weight="bold", labelpad=10)
plt.ylabel("Business Category", fontsize=12, weight="bold", labelpad=10)
plt.title("Average Security Score by Industry Category", fontsize=14, weight="bold", pad=15, color="#E5A88B")
plt.xlim(0, 5.5)
ax.grid(True, axis='x', linestyle='--', color='#312C45')
ax.tick_params(labelsize=11.5)
for bar in bars:
    width = bar.get_width()
    ax.text(width + 0.1, bar.get_y() + bar.get_height()/2, f'{width:.2f}', 
             va='center', ha='left', fontsize=12, weight='bold', color='#FAF6F0')
plt.tight_layout()
plt.savefig("charts/category_scores.png", dpi=300, facecolor='#120F1D')
plt.close()
print("Generating location scores bar chart...")
loc_scores = df.groupby("location")["security_score"].mean().sort_values(ascending=False)
fig, ax = plt.subplots(figsize=(8, 4.5))
bars = ax.bar(loc_scores.index, loc_scores.values, color="#E5A88B", width=0.45, edgecolor='#312C45', alpha=0.95) 
plt.xlabel("Neighborhood / Area in Delhi", fontsize=12, weight="bold", labelpad=10)
plt.ylabel("Average Security Score (0 - 5)", fontsize=12, weight="bold", labelpad=10)
plt.title("Geographical Security Analysis (Average Score)", fontsize=14, weight="bold", pad=15, color="#E5A88B")
plt.ylim(0, 5.5)
ax.grid(True, axis='y', linestyle='--', color='#312C45')
ax.tick_params(labelsize=11.5)
for bar in bars:
    height = bar.get_height()
    ax.text(bar.get_x() + bar.get_width()/2., height + 0.1, f'{height:.2f}', 
             ha='center', va='bottom', fontsize=11, weight='bold', color='#FAF6F0')
plt.xticks(rotation=15)
plt.tight_layout()
plt.savefig("charts/location_risks.png", dpi=300, facecolor='#120F1D')
plt.close()
print("Generating risk indicators frequency chart...")
indicators_flat = [indicator for sublist in df["risk_indicators"] for indicator in sublist]
indicator_counts = pd.Series(indicators_flat).value_counts().sort_values(ascending=True)
fig, ax = plt.subplots(figsize=(8, 4.5))
bars = ax.barh(indicator_counts.index, indicator_counts.values, color="#D5C3C6", height=0.55, edgecolor='#312C45', alpha=0.95) 
plt.xlabel("Number of Businesses Affected", fontsize=12, weight="bold", labelpad=10)
plt.ylabel("Security Vulnerability / Risk Indicator", fontsize=12, weight="bold", labelpad=10)
plt.title("Prevalence of Digital Risk Indicators", fontsize=14, weight="bold", pad=15, color="#E5A88B")
plt.xlim(0, max(indicator_counts.values) + 1.8)
ax.grid(True, axis='x', linestyle='--', color='#312C45')
ax.tick_params(labelsize=11.5)
for bar in bars:
    width = bar.get_width()
    ax.text(width + 0.15, bar.get_y() + bar.get_height()/2, f'{int(width)}', 
             va='center', ha='left', fontsize=12, weight='bold', color='#FAF6F0')
plt.tight_layout()
plt.savefig("charts/risk_indicators.png", dpi=300, facecolor='#120F1D')
plt.close()
print("Charts saved successfully in /charts folder.")
print("Generating Excel...")
wb = Workbook()
ws_dash = wb.active
ws_dash.title = "Executive Dashboard"
ws_data = wb.create_sheet(title="Security Assessment Data")
font_title = Font(name="Segoe UI", size=16, bold=True, color="231F33")
font_section = Font(name="Segoe UI", size=12, bold=True, color="231F33")
font_header = Font(name="Segoe UI", size=11, bold=True, color="FFFFFF")
font_body = Font(name="Segoe UI", size=10, color="2E2E2E")
font_body_bold = Font(name="Segoe UI", size=10, bold=True, color="2E2E2E")
font_kpi_num = Font(name="Segoe UI", size=20, bold=True, color="231F33")
font_kpi_label = Font(name="Segoe UI", size=9, color="7A7075", italic=True)
fill_header = PatternFill(start_color="231F33", end_color="231F33", fill_type="solid") 
fill_zebra = PatternFill(start_color="FAF7FA", end_color="FAF7FA", fill_type="solid") 
fill_kpi = PatternFill(start_color="F5EFEF", end_color="F5EFEF", fill_type="solid") 
fill_section_hdr = PatternFill(start_color="F9EBE5", end_color="F9EBE5", fill_type="solid") 
fills_risk = {
    "Very Secure": PatternFill(start_color="E8F8F0", end_color="E8F8F0", fill_type="solid"),
    "Secure": PatternFill(start_color="EDF7ED", end_color="EDF7ED", fill_type="solid"),
    "Moderate Risk": PatternFill(start_color="FFF9E6", end_color="FFF9E6", fill_type="solid"),
    "High Risk": PatternFill(start_color="FDF0E9", end_color="FDF0E9", fill_type="solid"),
    "Very High Risk": PatternFill(start_color="FDF1F0", end_color="FDF1F0", fill_type="solid")
}
fonts_risk = {
    "Very Secure": Font(name="Segoe UI", size=10, bold=True, color="1E6B43"),
    "Secure": Font(name="Segoe UI", size=10, bold=True, color="2E7D32"),
    "Moderate Risk": Font(name="Segoe UI", size=10, bold=True, color="B78103"),
    "High Risk": Font(name="Segoe UI", size=10, bold=True, color="A0522D"),
    "Very High Risk": Font(name="Segoe UI", size=10, bold=True, color="C0392B")
}
fills_priority = {
    "CRITICAL": PatternFill(start_color="FDF1F0", end_color="FDF1F0", fill_type="solid"),
    "MEDIUM": PatternFill(start_color="FFF9E6", end_color="FFF9E6", fill_type="solid"),
    "LOW": PatternFill(start_color="E8F8F0", end_color="E8F8F0", fill_type="solid")
}
fonts_priority = {
    "CRITICAL": Font(name="Segoe UI", size=10, bold=True, color="C0392B"),
    "MEDIUM": Font(name="Segoe UI", size=10, bold=True, color="B78103"),
    "LOW": Font(name="Segoe UI", size=10, bold=True, color="1E6B43")
}
thin_line = Side(style='thin', color='D9D9D9')
double_line = Side(style='double', color='000000')
border_all = Border(left=thin_line, right=thin_line, top=thin_line, bottom=thin_line)
border_header = Border(left=thin_line, right=thin_line, top=thin_line, bottom=Side(style='medium', color='1F4E78'))
border_total = Border(top=thin_line, bottom=double_line)
align_left = Alignment(horizontal='left', vertical='center')
align_center = Alignment(horizontal='center', vertical='center')
align_right = Alignment(horizontal='right', vertical='center')
align_left_wrap = Alignment(horizontal='left', vertical='center', wrap_text=True)
def create_dashboard_kpi(ws, col_start, col_end, formula_val, label_text, is_float=False):
    ws.merge_cells(f"{col_start}5:{col_end}5")
    ws.merge_cells(f"{col_start}6:{col_end}6")
    num_cell = ws[f"{col_start}5"]
    num_cell.value = formula_val
    num_cell.font = font_kpi_num
    num_cell.alignment = align_center
    num_cell.fill = fill_kpi
    if is_float:
        num_cell.number_format = '0.00'
    else:
        num_cell.number_format = '#,##0'
    lbl_cell = ws[f"{col_start}6"]
    lbl_cell.value = label_text
    lbl_cell.font = font_kpi_label
    lbl_cell.alignment = align_center
    lbl_cell.fill = fill_kpi
    for r in [5, 6]:
        col_ord_start = ord(col_start[0])
        col_ord_end = ord(col_end[0])
        for c_ord in range(col_ord_start, col_ord_end + 1):
            cell = ws[f"{chr(c_ord)}{r}"]
            cell.border = border_all
ws_data.views.sheetView[0].showGridLines = True
headers = [
    "Business ID", "Business Name", "Category", "Location", "Google Maps Rating", 
    "Reviews Count", "GBP Claimed", "Website Status", "Website URL", "HTTPS Redirect", 
    "SSL Status", "SSL Provider", "Admin Panel Exposed", "Forms Protection", 
    "Directory Listing", "Safe Browsing Check", "Missing Security Headers", "Security Score", 
    "Risk Level", "Outreach Priority", "Estimated Lead Value", "Outreach Strategy", 
    "Outreach Subject Line", "First Touch Pitch Hook", "Sales Pitch Hook", 
    "Key Security Gaps", "Actionable Recommendation"
]
for col_idx, h in enumerate(headers, 1):
    cell = ws_data.cell(row=2, column=col_idx, value=h)
    cell.font = font_header
    cell.fill = fill_header
    cell.alignment = align_center
    cell.border = border_header
ws_data.row_dimensions[2].height = 30
data_start_row = 3
for r_idx, d in enumerate(processed_data, data_start_row):
    ws_data.cell(row=r_idx, column=1, value=d["id"]).alignment = align_center
    ws_data.cell(row=r_idx, column=2, value=d["name"]).alignment = align_left
    ws_data.cell(row=r_idx, column=3, value=d["category"]).alignment = align_left
    ws_data.cell(row=r_idx, column=4, value=d["location"]).alignment = align_center
    r_cell = ws_data.cell(row=r_idx, column=5, value=d["rating"])
    r_cell.alignment = align_right
    r_cell.number_format = '0.0'
    rev_cell = ws_data.cell(row=r_idx, column=6, value=d["reviews"])
    rev_cell.alignment = align_right
    rev_cell.number_format = '#,##0'
    ws_data.cell(row=r_idx, column=7, value=d["gbp_claimed"]).alignment = align_center
    ws_data.cell(row=r_idx, column=8, value=d["website_status"]).alignment = align_center
    ws_data.cell(row=r_idx, column=9, value=d["website_url"]).alignment = align_left
    ws_data.cell(row=r_idx, column=10, value=d["https_redirect"]).alignment = align_center
    ws_data.cell(row=r_idx, column=11, value=d["ssl_status"]).alignment = align_center
    ws_data.cell(row=r_idx, column=12, value=d["ssl_provider"]).alignment = align_center
    ws_data.cell(row=r_idx, column=13, value=d["admin_panel_exposed"]).alignment = align_center
    ws_data.cell(row=r_idx, column=14, value=d["forms_protection"]).alignment = align_center
    ws_data.cell(row=r_idx, column=15, value=d["directory_listing"]).alignment = align_center
    ws_data.cell(row=r_idx, column=16, value=d["safe_browsing"]).alignment = align_center
    ws_data.cell(row=r_idx, column=17, value=d["missing_headers"]).alignment = align_left
    score_cell = ws_data.cell(row=r_idx, column=18, value=d["security_score"])
    score_cell.alignment = align_center
    score_cell.number_format = '0'
    risk_cell = ws_data.cell(row=r_idx, column=19, value=d["risk_level"])
    risk_cell.alignment = align_center
    risk_cell.fill = fills_risk[d["risk_level"]]
    risk_cell.font = fonts_risk[d["risk_level"]]
    priority_cell = ws_data.cell(row=r_idx, column=20, value=d["outreach_priority"])
    priority_cell.alignment = align_center
    priority_cell.fill = fills_priority.get(d["outreach_priority"])
    priority_cell.font = fonts_priority.get(d["outreach_priority"])
    ws_data.cell(row=r_idx, column=21, value=d["estimated_lead_value"]).alignment = align_center
    ws_data.cell(row=r_idx, column=22, value=d["outreach_strategy"]).alignment = align_left_wrap
    ws_data.cell(row=r_idx, column=23, value=d["outreach_subject_line"]).alignment = align_left_wrap
    ws_data.cell(row=r_idx, column=24, value=d["first_touch_pitch"]).alignment = align_left_wrap
    ws_data.cell(row=r_idx, column=25, value=d["sales_pitch_hook"]).alignment = align_left_wrap
    ws_data.cell(row=r_idx, column=26, value=d["key_security_gaps"]).alignment = align_left_wrap
    ws_data.cell(row=r_idx, column=27, value=d["recommendation"]).alignment = align_left_wrap
    for col_idx in range(1, len(headers) + 1):
        c = ws_data.cell(row=r_idx, column=col_idx)
        if col_idx not in [19, 20]:
            c.font = font_body
            if r_idx % 2 == 0:
                c.fill = fill_zebra
        c.border = border_all
    max_lines = 1
    for col_idx in range(1, len(headers) + 1):
        cell_val = str(ws_data.cell(row=r_idx, column=col_idx).value or '')
        col_width = 12 
        if col_idx == 23: col_width = 45 
        elif col_idx == 22: col_width = 30 
        elif col_idx == 21: col_width = 35 
        elif col_idx == 17: col_width = 25 
        elif col_idx == 9: col_width = 30 
        lines = (len(cell_val) // col_width) + 1
        if lines > max_lines:
            max_lines = lines
    ws_data.row_dimensions[r_idx].height = max(18 * max_lines, 26)
data_end_row = data_start_row + len(processed_data) - 1
tot_row = data_end_row + 2
ws_data.cell(row=tot_row, column=1, value="Summary Metrics").font = font_body_bold
ws_data.cell(row=tot_row, column=1).alignment = align_left
ws_data.cell(row=tot_row, column=5, value=f"=AVERAGE(E3:E{data_end_row})").font = font_body_bold
ws_data.cell(row=tot_row, column=5).alignment = align_right
ws_data.cell(row=tot_row, column=5).number_format = '0.00'
ws_data.cell(row=tot_row, column=6, value=f"=SUM(F3:F{data_end_row})").font = font_body_bold
ws_data.cell(row=tot_row, column=6).alignment = align_right
ws_data.cell(row=tot_row, column=6).number_format = '#,##0'
ws_data.cell(row=tot_row, column=18, value=f"=AVERAGE(R3:R{data_end_row})").font = font_body_bold
ws_data.cell(row=tot_row, column=18).alignment = align_center
ws_data.cell(row=tot_row, column=18).number_format = '0.00'
for col_idx in range(1, len(headers) + 1):
    c = ws_data.cell(row=tot_row, column=col_idx)
    c.border = border_total
ws_data.row_dimensions[tot_row].height = 24
col_widths = {
    1: 12,  
    2: 25,  
    3: 22,  
    4: 15,  
    5: 15,  
    6: 14,  
    7: 14,  
    8: 15,  
    9: 25,  
    10: 15, 
    11: 15, 
    12: 18, 
    13: 18, 
    14: 18, 
    15: 18, 
    16: 18, 
    17: 22, 
    18: 14, 
    19: 16, 
    20: 16, 
    21: 20, 
    22: 25, 
    23: 30, 
    24: 35, 
    25: 35, 
    26: 30, 
    27: 40, 
}
for col_idx, width in col_widths.items():
    col_letter = get_column_letter(col_idx)
    ws_data.column_dimensions[col_letter].width = width
ws_dash.views.sheetView[0].showGridLines = True
ws_dash.merge_cells("A2:H2")
ws_dash["A2"] = "BUSINESS SECURITY & DIGITAL RISK ANALYSIS - PAN-INDIA"
ws_dash["A2"].font = Font(name="Segoe UI", size=18, bold=True, color="FFFFFF")
ws_dash["A2"].fill = PatternFill(start_color="0F2027", end_color="0F2027", fill_type="solid")
ws_dash["A2"].alignment = Alignment(horizontal='center', vertical='center')
ws_dash.row_dimensions[2].height = 40
ws_dash.merge_cells("A3:H3")
ws_dash["A3"] = "Digital Security Posture Assessment of 50 Businesses Across India"
ws_dash["A3"].font = Font(name="Segoe UI", size=11, italic=True, color="CCCCCC")
ws_dash["A3"].fill = PatternFill(start_color="203A43", end_color="203A43", fill_type="solid")
ws_dash["A3"].alignment = Alignment(horizontal='center', vertical='center')
ws_dash.row_dimensions[3].height = 22
risk_rows = [
    ("Very Secure", f"=COUNTIF('Security Assessment Data'!S3:S{data_end_row}, \"Very Secure\")"),
    ("Secure", f"=COUNTIF('Security Assessment Data'!S3:S{data_end_row}, \"Secure\")"),
    ("Moderate Risk", f"=COUNTIF('Security Assessment Data'!S3:S{data_end_row}, \"Moderate Risk\")"),
    ("High Risk", f"=COUNTIF('Security Assessment Data'!S3:S{data_end_row}, \"High Risk\")"),
    ("Very High Risk", f"=COUNTIF('Security Assessment Data'!S3:S{data_end_row}, \"Very High Risk\")")
]
cat_rows = [
    "Restaurants & Cafes",
    "Healthcare & Clinics",
    "Retail & Shops",
    "Professional Services",
    "Hospitality & Hotels"
]
total_checked = len(processed_data)
no_website_count = sum(1 for d in processed_data if d["website_status"] == "No Website")
exposed_admin_count = sum(1 for d in processed_data if d["admin_panel_exposed"] == "Yes")
unclaimed_count = sum(1 for d in processed_data if d["gbp_claimed"] == "No")
active_site_count = total_checked - no_website_count
missing_headers_count = sum(1 for d in processed_data if d["website_status"] == "Active" and d["missing_headers"] != "None")
missing_headers_pct = (missing_headers_count / active_site_count) * 100 if active_site_count > 0 else 0
observations = [
    "1. High Vulnerability in Professional Services: Professional services have a low average security score. Several regional law offices and tax consultants exhibit critical weaknesses.",
    f"2. Digital Absence Risk: {no_website_count} out of {total_checked} businesses analyzed have NO website, leaving them with a Very High Risk rating due to missed opportunities.",
    f"3. Missing Security Headers: {missing_headers_pct:.0f}% of the websites scanned are missing critical security headers like Content-Security-Policy (CSP) and HSTS.",
    f"4. Exposed Administrative Interfaces: {exposed_admin_count} businesses have exposed admin panels (/wp-admin or similar) accessible to the public.",
    f"5. Unclaimed Profiles: {unclaimed_count} businesses have unclaimed Google Business Profiles on Maps. This exposes them to direct profile hijacking and sabotage."
]
create_dashboard_kpi(ws_dash, "A", "B", f"=COUNTA('Security Assessment Data'!A3:A{data_end_row})", "Total Businesses Checked")
create_dashboard_kpi(ws_dash, "C", "D", f"=AVERAGE('Security Assessment Data'!R3:R{data_end_row})", "Average Security Score (0-5)", is_float=True)
create_dashboard_kpi(ws_dash, "E", "F", f"=COUNTIF('Security Assessment Data'!S3:S{data_end_row}, \"*Secure*\")", "Secure Businesses (Score 4-5)")
create_dashboard_kpi(ws_dash, "G", "H", f"=COUNTIF('Security Assessment Data'!S3:S{data_end_row}, \"*Risk*\")", "Risky Businesses (Score 0-3)")
ws_dash.row_dimensions[5].height = 32
ws_dash.row_dimensions[6].height = 18
ws_dash.row_dimensions[9].height = 24
ws_dash["A9"] = "Risk Level Breakdown"
ws_dash["A9"].font = font_section
ws_dash["A9"].fill = fill_section_hdr
ws_dash["A9"].alignment = align_left
ws_dash.merge_cells("A9:C9")
ws_dash["A9"].border = border_all
ws_dash["B9"].border = border_all
ws_dash["C9"].border = border_all
risk_headers = ["Risk Level", "Count", "Percentage"]
for idx, h in enumerate(risk_headers, 1):
    c = ws_dash.cell(row=10, column=idx, value=h)
    c.font = font_header
    c.fill = fill_header
    c.alignment = align_center
    c.border = border_header
ws_dash.row_dimensions[10].height = 24
for idx, (level, formula) in enumerate(risk_rows, 11):
    c_lvl = ws_dash.cell(row=idx, column=1, value=level)
    c_lvl.font = font_body
    c_lvl.border = border_all
    c_lvl.alignment = align_left
    c_cnt = ws_dash.cell(row=idx, column=2, value=formula)
    c_cnt.font = font_body
    c_cnt.border = border_all
    c_cnt.alignment = align_center
    c_pct = ws_dash.cell(row=idx, column=3, value=f"=B{idx}/SUM(B11:B15)")
    c_pct.font = font_body
    c_pct.border = border_all
    c_pct.alignment = align_right
    c_pct.number_format = '0.0%'
    ws_dash.row_dimensions[idx].height = 20
tot_level_row = 16
ws_dash.cell(row=tot_level_row, column=1, value="Total").font = font_body_bold
ws_dash.cell(row=tot_level_row, column=1).alignment = align_left
ws_dash.cell(row=tot_level_row, column=2, value="=SUM(B11:B15)").font = font_body_bold
ws_dash.cell(row=tot_level_row, column=2).alignment = align_center
ws_dash.cell(row=tot_level_row, column=3, value="=SUM(C11:C15)").font = font_body_bold
ws_dash.cell(row=tot_level_row, column=3).alignment = align_right
ws_dash.cell(row=tot_level_row, column=3).number_format = '0.0%'
for c_idx in range(1, 4):
    ws_dash.cell(row=tot_level_row, column=c_idx).border = border_total
ws_dash.column_dimensions['D'].width = 4
ws_dash["E9"] = "Average Score by Category"
ws_dash["E9"].font = font_section
ws_dash["E9"].fill = fill_section_hdr
ws_dash["E9"].alignment = align_left
ws_dash.merge_cells("E9:G9")
ws_dash["E9"].border = border_all
ws_dash["F9"].border = border_all
ws_dash["G9"].border = border_all
cat_headers = ["Category", "Total Checked", "Average Score"]
for idx, h in enumerate(cat_headers, 5):
    c = ws_dash.cell(row=10, column=idx, value=h)
    c.font = font_header
    c.fill = fill_header
    c.alignment = align_center
    c.border = border_header
for idx, cat in enumerate(cat_rows, 11):
    c_cat = ws_dash.cell(row=idx, column=5, value=cat)
    c_cat.font = font_body
    c_cat.border = border_all
    c_cat.alignment = align_left
    c_cnt = ws_dash.cell(row=idx, column=6, value=f"=COUNTIF('Security Assessment Data'!C3:C{data_end_row}, \"{cat}\")")
    c_cnt.font = font_body
    c_cnt.border = border_all
    c_cnt.alignment = align_center
    c_avg = ws_dash.cell(row=idx, column=7, value=f"=AVERAGEIF('Security Assessment Data'!C3:C{data_end_row}, \"{cat}\", 'Security Assessment Data'!R3:R{data_end_row})")
    c_avg.font = font_body
    c_avg.border = border_all
    c_avg.alignment = align_center
    c_avg.number_format = '0.00'
    ws_dash.row_dimensions[idx].height = 20
tot_cat_row = 16
ws_dash.cell(row=tot_cat_row, column=5, value="All Industries").font = font_body_bold
ws_dash.cell(row=tot_cat_row, column=5).alignment = align_left
ws_dash.cell(row=tot_cat_row, column=6, value="=SUM(F11:F15)").font = font_body_bold
ws_dash.cell(row=tot_cat_row, column=6).alignment = align_center
ws_dash.cell(row=tot_cat_row, column=7, value="=AVERAGE('Security Assessment Data'!R3:R{data_end_row})").font = font_body_bold
ws_dash.cell(row=tot_cat_row, column=7).alignment = align_center
ws_dash.cell(row=tot_cat_row, column=7).number_format = '0.00'
for c_idx in range(5, 8):
    ws_dash.cell(row=tot_cat_row, column=c_idx).border = border_total
ws_dash["A19"] = "Key Observation & Analysis Insights:"
ws_dash["A19"].font = font_section
ws_dash["A19"].fill = fill_section_hdr
ws_dash.merge_cells("A19:H19")
for c_ord in range(ord('A'), ord('H') + 1):
    ws_dash[f"{chr(c_ord)}19"].border = border_all
ws_dash.row_dimensions[19].height = 24
for idx, obs in enumerate(observations, 20):
    ws_dash.cell(row=idx, column=1, value=obs).font = Font(name="Segoe UI", size=10, color="555555")
    ws_dash.cell(row=idx, column=1).alignment = align_left
    ws_dash.merge_cells(start_row=idx, start_column=1, end_row=idx, end_column=8)
    ws_dash.row_dimensions[idx].height = 20
ws_dash.column_dimensions['A'].width = 24
ws_dash.column_dimensions['B'].width = 12
ws_dash.column_dimensions['C'].width = 12
ws_dash.column_dimensions['D'].width = 4
ws_dash.column_dimensions['E'].width = 26
ws_dash.column_dimensions['F'].width = 15
ws_dash.column_dimensions['G'].width = 15
ws_dash.column_dimensions['H'].width = 12
wb.save("National_Digital_Risk_Audit_Database.xlsx")
print("Excel Workbook National_Digital_Risk_Audit_Database.xlsx generated and styled successfully!")
print("Generating PPTX presentation...")
prs = Presentation()
prs.slide_width = Inches(16.0) 
prs.slide_height = Inches(9.0)
DARK_BG = RGBColor(18, 15, 29)        
CARD_BG = RGBColor(27, 24, 41)        
BORDER_COLOR = RGBColor(49, 44, 69)   
TEXT_WHITE = RGBColor(250, 246, 240)  
TEXT_MUTED = RGBColor(213, 195, 198)  
CYAN_ACCENT = RGBColor(229, 168, 139)  
GREEN_SECURE = RGBColor(136, 212, 158)  
RED_RISK = RGBColor(240, 113, 103)    
YELLOW_WARN = RGBColor(241, 196, 15)   
def shorten_gaps(gaps):
    gaps = gaps.replace("Exposed Admin Panel (/wp-admin or /admin)", "Exposed Admin Page")
    gaps = gaps.replace("Digital Absence, Unclaimed Google Business Profile", "No Website, Unclaimed Profile")
    gaps = gaps.replace("Digital Absence, No Website Present", "No Website (Digital Absence)")
    gaps = gaps.replace("No HTTPS Redirect", "HTTP Only (No HTTPS)")
    gaps = gaps.replace("Unprotected Forms, Directory Listing Enabled", "Unprotected Forms, Open Dir")
    gaps = gaps.replace("All Headers Missing", "Missing Headers")
    gaps = gaps.replace("SSL Status: Expired", "Expired SSL")
    gaps = gaps.replace("SSL Status: None", "No SSL")
    return gaps
def apply_background(slide, color=DARK_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color
def create_title(slide, text, subtitle_text=None, top_inches=0.5):
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(top_inches), Inches(14.0), Inches(1.1))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Segoe UI"
    p.font.size = Pt(44) 
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(18) 
        p2.font.color.rgb = CYAN_ACCENT
        p2.space_before = Pt(4)
    badge_box = slide.shapes.add_textbox(Inches(11.0), Inches(top_inches), Inches(4.0), Inches(0.5))
    tf_b = badge_box.text_frame
    tf_b.word_wrap = False
    tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
    p_b = tf_b.paragraphs[0]
    p_b.text = "PAN-INDIA DIGITAL SECURITY AUDIT • 2026"
    p_b.alignment = PP_ALIGN.RIGHT
    p_b.font.name = "Segoe UI"
    p_b.font.size = Pt(10)
    p_b.font.bold = True
    p_b.font.color.rgb = TEXT_MUTED
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(1.0), Inches(top_inches + 1.2), Inches(14.0), Inches(0.05)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = CYAN_ACCENT
    divider.line.color.rgb = CYAN_ACCENT
    divider.line.width = Pt(1)
def draw_premium_card(slide, x, y, w, h, bg_color=CARD_BG, border_color=BORDER_COLOR, accent_color=CYAN_ACCENT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.04
    if accent_color:
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(0.08)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = accent_color
        accent.line.color.rgb = accent_color
        accent.line.width = Pt(1)
    return shape
slide_layout = prs.slide_layouts[6] 
slide1 = prs.slides.add_slide(slide_layout)
apply_background(slide1)
draw_premium_card(slide1, 1.0, 2.0, 14.0, 5.0, bg_color=CARD_BG, border_color=BORDER_COLOR, accent_color=CYAN_ACCENT)
tb = slide1.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(13.0), Inches(4.0))
tf = tb.text_frame
tf.word_wrap = True
p_sub = tf.paragraphs[0]
p_sub.text = "WEEK 3 - CYBER SECURITY TASK"
p_sub.font.name = "Segoe UI"
p_sub.font.size = Pt(20) 
p_sub.font.bold = True
p_sub.font.color.rgb = CYAN_ACCENT
p_main = tf.add_paragraph()
p_main.text = "BUSINESS SECURITY & DIGITAL RISK ANALYSIS"
p_main.font.name = "Segoe UI"
p_main.font.size = Pt(44) 
p_main.font.bold = True
p_main.font.color.rgb = TEXT_WHITE
p_main.space_before = Pt(12)
p_desc = tf.add_paragraph()
p_desc.text = "Digital presence assessment, vulnerability scan, and cyber risk scoring of 50 businesses identified via Google Maps across major Indian cities."
p_desc.font.name = "Segoe UI"
p_desc.font.size = Pt(16) 
p_desc.font.color.rgb = TEXT_MUTED
p_desc.space_before = Pt(20)
p_author = tf.add_paragraph()
p_author.text = "Prepared by: Cyber Security Consultant  |  Assessment Scope: Pan-India Market Audit"
p_author.font.name = "Segoe UI"
p_author.font.size = Pt(14) 
p_author.font.color.rgb = CYAN_ACCENT
p_author.space_before = Pt(50)
slide2 = prs.slides.add_slide(slide_layout)
apply_background(slide2)
create_title(slide2, "Project Scope & Assessment Areas", "Target Region: India Local Market Analysis (50 Businesses)")
col_w = 4.4
gap = 0.4
y_pos = 2.4
h_pos = 5.6
draw_premium_card(slide2, 1.0, y_pos, col_w, h_pos, accent_color=CYAN_ACCENT)
tb_sc1 = slide2.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.3), Inches(col_w - 0.4), Inches(h_pos - 0.6))
tf = tb_sc1.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "01  |  Data Sourcing"
p.font.name = "Segoe UI"
p.font.size = Pt(26) 
p.font.bold = True
p.font.color.rgb = CYAN_ACCENT
p.space_after = Pt(16)
p_body1_hdr = tf.add_paragraph()
p_body1_hdr.text = "Google Maps Sourcing"
p_body1_hdr.font.name = "Segoe UI"
p_body1_hdr.font.size = Pt(20) 
p_body1_hdr.font.bold = True
p_body1_hdr.font.color.rgb = TEXT_WHITE
p_body1 = tf.add_paragraph()
p_body1.text = "• Mapped 50 Indian businesses.\n• Mumbai, Bangalore, Chennai, Kolkata.\n• Delhi, Hyderabad, Pune, Jaipur, Goa."
p_body1.font.name = "Segoe UI"
p_body1.font.size = Pt(18) 
p_body1.font.color.rgb = TEXT_MUTED
p_body1.space_before = Pt(4)
p_body1.space_after = Pt(16)
p_body2_hdr = tf.add_paragraph()
p_body2_hdr.text = "Security Audits"
p_body2_hdr.font.name = "Segoe UI"
p_body2_hdr.font.size = Pt(20) 
p_body2_hdr.font.bold = True
p_body2_hdr.font.color.rgb = TEXT_WHITE
p_body2 = tf.add_paragraph()
p_body2.text = "• Scanned target websites.\n• Checked SSL & admin page flags."
p_body2.font.name = "Segoe UI"
p_body2.font.size = Pt(18) 
p_body2.font.color.rgb = TEXT_MUTED
p_body2.space_before = Pt(4)
draw_premium_card(slide2, 1.0 + col_w + gap, y_pos, col_w, h_pos, accent_color=CYAN_ACCENT)
tb_sc2 = slide2.shapes.add_textbox(Inches(1.2 + col_w + gap), Inches(y_pos + 0.3), Inches(col_w - 0.4), Inches(h_pos - 0.6))
tf = tb_sc2.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "02  |  Industry Scope"
p.font.name = "Segoe UI"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = CYAN_ACCENT
p.space_after = Pt(16)
p_body1_hdr = tf.add_paragraph()
p_body1_hdr.text = "Sectors Evaluated"
p_body1_hdr.font.name = "Segoe UI"
p_body1_hdr.font.size = Pt(20)
p_body1_hdr.font.bold = True
p_body1_hdr.font.color.rgb = TEXT_WHITE
p_body1 = tf.add_paragraph()
p_body1.text = "• Restaurants & Cafes (10 checked)\n• Healthcare & Clinics (10 checked)\n• Retail & Shops (10 checked)\n• Professional Services (10 checked)\n• Hospitality & Hotels (10 checked)"
p_body1.font.name = "Segoe UI"
p_body1.font.size = Pt(18) 
p_body1.font.color.rgb = TEXT_MUTED
p_body1.space_before = Pt(4)
p_body1.space_after = Pt(16)
draw_premium_card(slide2, 1.0 + 2*(col_w + gap), y_pos, col_w, h_pos, accent_color=CYAN_ACCENT)
tb_sc3 = slide2.shapes.add_textbox(Inches(1.2 + 2*(col_w + gap)), Inches(y_pos + 0.3), Inches(col_w - 0.4), Inches(h_pos - 0.6))
tf = tb_sc3.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "03  |  Security Parameters"
p.font.name = "Segoe UI"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = CYAN_ACCENT
p.space_after = Pt(16)
p_body1_hdr = tf.add_paragraph()
p_body1_hdr.text = "Protocols & SSL Checks"
p_body1_hdr.font.name = "Segoe UI"
p_body1_hdr.font.size = Pt(20)
p_body1_hdr.font.bold = True
p_body1_hdr.font.color.rgb = TEXT_WHITE
p_body1 = tf.add_paragraph()
p_body1.text = "• Checked HTTPS redirect.\n• Checked SSL certificate validity.\n• Scanned admin login exposure."
p_body1.font.name = "Segoe UI"
p_body1.font.size = Pt(18) 
p_body1.font.color.rgb = TEXT_MUTED
p_body1.space_before = Pt(4)
p_body1.space_after = Pt(16)
p_body2_hdr = tf.add_paragraph()
p_body2_hdr.text = "Asset Exposures"
p_body2_hdr.font.name = "Segoe UI"
p_body2_hdr.font.size = Pt(20)
p_body2_hdr.font.bold = True
p_body2_hdr.font.color.rgb = TEXT_WHITE
p_body2 = tf.add_paragraph()
p_body2.text = "• Scanned open directory indexes.\n• Tested contact forms protection."
p_body2.font.name = "Segoe UI"
p_body2.font.size = Pt(18) 
p_body2.font.color.rgb = TEXT_MUTED
p_body2.space_before = Pt(4)
slide3 = prs.slides.add_slide(slide_layout)
apply_background(slide3)
create_title(slide3, "Executive Summary", "Key Findings and Performance Metrics")
kpi_w = 3.2
kpi_gap = 0.4
y_kpi = 2.4
h_kpi = 1.8
total_checked = len(df)
avg_score = df["security_score"].mean()
secure_count = len(df[df["security_score"] >= 4])
secure_rate = (secure_count / total_checked) * 100
risk_rate = 100 - secure_rate
no_website_count = len(df[df["website_status"] == "No Website"])
unclaimed_count = len(df[df["gbp_claimed"] == "No"])
exposed_admin_count = len(df[df["admin_panel_exposed"] == "Yes"])
kpi_data = [
    {"num": str(total_checked), "lbl": "TOTAL CHECKED", "color": TEXT_WHITE, "accent": TEXT_MUTED},
    {"num": f"{avg_score:.2f}", "lbl": "AVERAGE SCORE", "color": CYAN_ACCENT, "accent": CYAN_ACCENT},
    {"num": f"{secure_rate:.0f}%", "lbl": "SECURE RATE", "color": GREEN_SECURE, "accent": GREEN_SECURE},
    {"num": f"{risk_rate:.0f}%", "lbl": "AT RISK RATE", "color": RED_RISK, "accent": RED_RISK}
]
for idx, kd in enumerate(kpi_data):
    x_k = 1.0 + idx * (kpi_w + kpi_gap)
    draw_premium_card(slide3, x_k, y_kpi, kpi_w, h_kpi, accent_color=kd["accent"])
    tb_k = slide3.shapes.add_textbox(Inches(x_k + 0.1), Inches(y_kpi + 0.1), Inches(kpi_w - 0.2), Inches(h_kpi - 0.2))
    tf = tb_k.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_num = tf.paragraphs[0]
    p_num.text = kd["num"]
    p_num.alignment = PP_ALIGN.CENTER
    p_num.font.name = "Segoe UI"
    p_num.font.size = Pt(48) 
    p_num.font.bold = True
    p_num.font.color.rgb = kd["color"]
    p_lbl = tf.add_paragraph()
    p_lbl.text = kd["lbl"]
    p_lbl.alignment = PP_ALIGN.CENTER
    p_lbl.font.name = "Segoe UI"
    p_lbl.font.size = Pt(15) 
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = TEXT_MUTED
    p_lbl.space_before = Pt(3)
draw_premium_card(slide3, 1.0, 4.6, 14.0, 3.4, accent_color=CYAN_ACCENT)
tb_summary = slide3.shapes.add_textbox(Inches(1.3), Inches(4.8), Inches(13.4), Inches(3.0))
tf = tb_summary.text_frame
tf.word_wrap = True
p_title = tf.paragraphs[0]
p_title.text = "High-Level Cybersecurity Posture Analysis"
p_title.font.name = "Segoe UI"
p_title.font.size = Pt(26) 
p_title.font.bold = True
p_title.font.color.rgb = CYAN_ACCENT
p_title.space_after = Pt(10)
p_text = tf.add_paragraph()
p_text.text = (
    f"• Security Divide: {secure_rate:.0f}% Secure rate is driven by corporate hotels/hospitals, but {risk_rate:.0f}% At-Risk rate exposes local small businesses.\n"
    f"• Protocol Gaps: Plaintext transmissions and exposed CMS panels ({exposed_admin_count} sites) are the largest source of website vulnerability.\n"
    f"• Sabotage Risks: {unclaimed_count} Google Maps business profiles remain unclaimed, presenting an active hijacking risk for information edits."
)
p_text.font.name = "Segoe UI"
p_text.font.size = Pt(18) 
p_text.font.color.rgb = TEXT_WHITE
p_text.space_before = Pt(4)
slide4 = prs.slides.add_slide(slide_layout)
apply_background(slide4)
create_title(slide4, "Security Risk & Score Distribution", "Proportion of overall secure vs. vulnerable businesses")
draw_premium_card(slide4, 1.0, 2.4, 6.5, 5.6, accent_color=CYAN_ACCENT)
tb_dist = slide4.shapes.add_textbox(Inches(1.3), Inches(2.6), Inches(5.9), Inches(5.2))
tf = tb_dist.text_frame
tf.word_wrap = True
p_h = tf.paragraphs[0]
p_h.text = "Risk Tier Analysis"
p_h.font.name = "Segoe UI"
p_h.font.size = Pt(26) 
p_h.font.bold = True
p_h.font.color.rgb = CYAN_ACCENT
p_h.space_after = Pt(14)
p1_title = tf.add_paragraph()
p1_title.text = f"Very Secure (Score 5) — {len(df[df['security_score'] == 5])} Businesses"
p1_title.font.name = "Segoe UI"
p1_title.font.bold = True
p1_title.font.size = Pt(20) 
p1_title.font.color.rgb = GREEN_SECURE
p1_desc = tf.add_paragraph()
p1_desc.text = "Fully secure websites with active HTTPS redirections, valid SSL certificates, and hidden CMS login pages."
p1_desc.font.name = "Segoe UI"
p1_desc.font.size = Pt(18) 
p1_desc.font.color.rgb = TEXT_WHITE
p1_desc.space_after = Pt(12)
p2_title = tf.add_paragraph()
p2_title.text = f"Secure (Score 4) — {len(df[df['security_score'] == 4])} Businesses"
p2_title.font.name = "Segoe UI"
p2_title.font.bold = True
p2_title.font.size = Pt(20) 
p2_title.font.color.rgb = CYAN_ACCENT
p2_desc = tf.add_paragraph()
p2_desc.text = "Minor vulnerabilities detected (missing CSP or HSTS headers), but user input forms are secure."
p2_desc.font.name = "Segoe UI"
p2_desc.font.size = Pt(18) 
p2_desc.font.color.rgb = TEXT_WHITE
p2_desc.space_after = Pt(12)
p3_title = tf.add_paragraph()
p3_title.text = f"Vulnerable (Score 0-3) — {len(df[df['security_score'] <= 3])} Businesses"
p3_title.font.name = "Segoe UI"
p3_title.font.bold = True
p3_title.font.size = Pt(20) 
p3_title.font.color.rgb = RED_RISK
p3_desc = tf.add_paragraph()
p3_desc.text = "Critical protocol exposures or complete digital absence (no website) with unclaimed business profiles."
p3_desc.font.name = "Segoe UI"
p3_desc.font.size = Pt(18) 
p3_desc.font.color.rgb = TEXT_WHITE
chart_path_1 = "charts/risk_distribution.png"
if os.path.exists(chart_path_1):
    slide4.shapes.add_picture(chart_path_1, Inches(8.0), Inches(2.4), width=Inches(7.0), height=Inches(5.6))
slide5 = prs.slides.add_slide(slide_layout)
apply_background(slide5)
create_title(slide5, "Vulnerability & Risk Indicators", "Prevalence of specific digital risks across the assessed group")
chart_path_2 = "charts/risk_indicators.png"
if os.path.exists(chart_path_2):
    slide5.shapes.add_picture(chart_path_2, Inches(1.0), Inches(2.4), width=Inches(7.0), height=Inches(5.6))
draw_premium_card(slide5, 8.5, 2.4, 6.5, 5.6, accent_color=RED_RISK)
tb_vul = slide5.shapes.add_textbox(Inches(8.8), Inches(2.6), Inches(5.9), Inches(5.2))
tf = tb_vul.text_frame
tf.word_wrap = True
p_h = tf.paragraphs[0]
p_h.text = "Vulnerability Breakdown"
p_h.font.name = "Segoe UI"
p_h.font.size = Pt(26) 
p_h.font.bold = True
p_h.font.color.rgb = CYAN_ACCENT
p_h.space_after = Pt(14)
p1_title = tf.add_paragraph()
insecure_count = len(df[(df["website_status"] == "Active") & ((df["ssl_status"] == "None") | (df["ssl_status"] == "Expired") | (df["https_redirect"] == "No"))])
p1_title.text = f"Insecure Protocols ({insecure_count} Sites)"
p1_title.font.name = "Segoe UI"
p1_title.font.bold = True
p1_title.font.size = Pt(20) 
p1_title.font.color.rgb = RED_RISK
p1_desc = tf.add_paragraph()
p1_desc.text = "Transmitting booking and user details in plaintext over HTTP. Vulnerable to interception."
p1_desc.font.name = "Segoe UI"
p1_desc.font.size = Pt(18) 
p1_desc.font.color.rgb = TEXT_WHITE
p1_desc.space_after = Pt(12)
p2_title = tf.add_paragraph()
exposed_admin_count = len(df[df["admin_panel_exposed"] == "Yes"])
p2_title.text = f"Exposed Admin Panels ({exposed_admin_count} Sites)"
p2_title.font.name = "Segoe UI"
p2_title.font.bold = True
p2_title.font.size = Pt(20) 
p2_title.font.color.rgb = RED_RISK
p2_desc = tf.add_paragraph()
p2_desc.text = "WordPress login directories are exposed, allowing hackers to run automated brute-force scans."
p2_desc.font.name = "Segoe UI"
p2_desc.font.size = Pt(18) 
p2_desc.font.color.rgb = TEXT_WHITE
p2_desc.space_after = Pt(12)
p3_title = tf.add_paragraph()
open_dir_count = len(df[df["directory_listing"] == "Yes"])
p3_title.text = f"Open Directories ({open_dir_count} Sites)"
p3_title.font.name = "Segoe UI"
p3_title.font.bold = True
p3_title.font.size = Pt(20) 
p3_title.font.color.rgb = RED_RISK
p3_desc = tf.add_paragraph()
p3_desc.text = "Folder indexing is active, exposing backend assets and file directories to public view."
p3_desc.font.name = "Segoe UI"
p3_desc.font.size = Pt(18) 
p3_desc.font.color.rgb = TEXT_WHITE
slide6 = prs.slides.add_slide(slide_layout)
apply_background(slide6)
create_title(slide6, "Industry Category & Location Analysis", "Comparing security standards across sectors and areas")
draw_premium_card(slide6, 1.0, 2.4, 6.5, 5.6, accent_color=CYAN_ACCENT)
tb_loc = slide6.shapes.add_textbox(Inches(1.3), Inches(2.6), Inches(5.9), Inches(5.2))
tf = tb_loc.text_frame
tf.word_wrap = True
p_h = tf.paragraphs[0]
p_h.text = "Industry & Geography Risks"
p_h.font.name = "Segoe UI"
p_h.font.size = Pt(26) 
p_h.font.bold = True
p_h.font.color.rgb = CYAN_ACCENT
p_h.space_after = Pt(14)
p1_title = tf.add_paragraph()
p1_title.text = "Industry Security Scorecard"
p1_title.font.name = "Segoe UI"
p1_title.font.bold = True
p1_title.font.size = Pt(20) 
p1_title.font.color.rgb = CYAN_ACCENT
p1_desc = tf.add_paragraph()
p1_desc.text = "Hospitality and Healthcare lead in security standards due to compliance regulations. Professional Services and Retail lag behind due to high rates of digital absence and unencrypted sites."
p1_desc.font.name = "Segoe UI"
p1_desc.font.size = Pt(18) 
p1_desc.font.color.rgb = TEXT_WHITE
p1_desc.space_after = Pt(14)
p2_title = tf.add_paragraph()
p2_title.text = "Geographical Analysis"
p2_title.font.name = "Segoe UI"
p2_title.font.bold = True
p2_title.font.size = Pt(20) 
p2_title.font.color.rgb = CYAN_ACCENT
p2_desc = tf.add_paragraph()
p2_desc.text = "Major metropolitan centers show stark contrasts: premium corporate entities in Mumbai and Bengaluru exhibit high security postures, while smaller regional firms represent high-vulnerability clusters."
p2_desc.font.name = "Segoe UI"
p2_desc.font.size = Pt(18) 
p2_desc.font.color.rgb = TEXT_WHITE
chart_path_3 = "charts/category_scores.png"
if os.path.exists(chart_path_3):
    slide6.shapes.add_picture(chart_path_3, Inches(8.0), Inches(2.4), width=Inches(7.0), height=Inches(5.6))
slide7 = prs.slides.add_slide(slide_layout)
apply_background(slide7)
create_title(slide7, "Top 10 Risky Businesses", "Identified targets requiring immediate intervention")
df_sorted = df.sort_values(by=["security_score", "reviews"], ascending=[True, False]).head(10)
rows, cols = 11, 5
left, top, width, height = Inches(1.0), Inches(2.4), Inches(14.0), Inches(5.8)
table_shape = slide7.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table
table.rows[0].height = Inches(0.60)
for r in range(1, 11):
    table.rows[r].height = Inches(0.65) 
table.columns[0].width = Inches(3.3)  
table.columns[1].width = Inches(2.4)  
table.columns[2].width = Inches(1.8)  
table.columns[3].width = Inches(1.5)  
table.columns[4].width = Inches(5.0) 
headers_tbl = ["Business Name", "Category", "Location", "Security Score", "Primary Security Gaps"]
for col_idx, text in enumerate(headers_tbl):
    cell = table.cell(0, col_idx)
    cell.text = text
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE 
    cell.fill.solid()
    cell.fill.fore_color.rgb = CARD_BG
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Segoe UI"
        p.font.size = Pt(16) 
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT
for row_idx, (_, r) in enumerate(df_sorted.iterrows(), 1):
    table.cell(row_idx, 0).text = r["name"]
    table.cell(row_idx, 1).text = r["category"]
    table.cell(row_idx, 2).text = r["location"]
    table.cell(row_idx, 3).text = f"{int(r['security_score'])} / 5"
    table.cell(row_idx, 4).text = shorten_gaps(r["key_security_gaps"])
    for col_idx in range(cols):
        cell = table.cell(row_idx, col_idx)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE 
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BG
        score_cell = table.cell(row_idx, 3)
        for p in score_cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            if r["security_score"] <= 1:
                p.font.color.rgb = RED_RISK
            elif r["security_score"] == 2:
                p.font.color.rgb = RGBColor(230, 126, 34)
            else:
                p.font.color.rgb = RGBColor(241, 196, 15)
        for p in cell.text_frame.paragraphs:
            p.font.name = "Segoe UI"
            p.font.size = Pt(14) 
            if col_idx != 3:
                p.font.color.rgb = TEXT_WHITE
            if col_idx == 4:
                p.alignment = PP_ALIGN.LEFT
            elif col_idx in [1, 2]:
                p.alignment = PP_ALIGN.CENTER
slide8 = prs.slides.add_slide(slide_layout)
apply_background(slide8)
create_title(slide8, "Strategic Cybersecurity Roadmap", "Actionable recommendations for businesses and service providers")
w_q, h_q = 6.8, 2.7
x_l, x_r = 1.0, 8.2
y_t, y_b = 2.4, 5.5
draw_premium_card(slide8, x_l, y_t, w_q, h_q, accent_color=RED_RISK)
tb_q1 = slide8.shapes.add_textbox(Inches(x_l + 0.2), Inches(y_t + 0.2), Inches(w_q - 0.4), Inches(h_q - 0.4))
tf = tb_q1.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "1. Immediate Remediation (At-Risk Sites)"
p.font.name = "Segoe UI"
p.font.size = Pt(20) 
p.font.bold = True
p.font.color.rgb = CYAN_ACCENT
p.space_after = Pt(8)
p_b1 = tf.add_paragraph()
p_b1.text = "• Install SSL Certificates: Force HTTP to HTTPS redirects."
p_b1.font.name = "Segoe UI"
p_b1.font.size = Pt(15) 
p_b1.font.color.rgb = TEXT_WHITE
p_b2 = tf.add_paragraph()
p_b2.text = "• Cloak Admin Panels: Restrict access to /wp-admin."
p_b2.font.name = "Segoe UI"
p_b2.font.size = Pt(15) 
p_b2.font.color.rgb = TEXT_WHITE
p_b3 = tf.add_paragraph()
p_b3.text = "• Secure Forms: Deploy anti-CSRF tokens on forms."
p_b3.font.name = "Segoe UI"
p_b3.font.size = Pt(15) 
p_b3.font.color.rgb = TEXT_WHITE
draw_premium_card(slide8, x_r, y_t, w_q, h_q, accent_color=CYAN_ACCENT)
tb_q2 = slide8.shapes.add_textbox(Inches(x_r + 0.2), Inches(y_t + 0.2), Inches(w_q - 0.4), Inches(h_q - 0.4))
tf = tb_q2.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "2. Google Maps Profile Hardening"
p.font.name = "Segoe UI"
p.font.size = Pt(20) 
p.font.bold = True
p.font.color.rgb = CYAN_ACCENT
p.space_after = Pt(8)
p_b1 = tf.add_paragraph()
p_b1.text = "• Profile Claiming: Verify unclaimed profiles on Maps."
p_b1.font.name = "Segoe UI"
p_b1.font.size = Pt(15) 
p_b1.font.color.rgb = TEXT_WHITE
p_b2 = tf.add_paragraph()
p_b2.text = "• Access Audits: Enforce MFA on associated accounts."
p_b2.font.name = "Segoe UI"
p_b2.font.size = Pt(15) 
p_b2.font.color.rgb = TEXT_WHITE
p_b3 = tf.add_paragraph()
p_b3.text = "• Monitor Edits: Set alerts for Google Maps suggestions."
p_b3.font.name = "Segoe UI"
p_b3.font.size = Pt(15) 
p_b3.font.color.rgb = TEXT_WHITE
draw_premium_card(slide8, x_l, y_b, w_q, h_q, accent_color=GREEN_SECURE)
tb_q3 = slide8.shapes.add_textbox(Inches(x_l + 0.2), Inches(y_b + 0.2), Inches(w_q - 0.4), Inches(h_q - 0.4))
tf = tb_q3.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "3. Ongoing Security Management"
p.font.name = "Segoe UI"
p.font.size = Pt(20) 
p.font.bold = True
p.font.color.rgb = CYAN_ACCENT
p.space_after = Pt(8)
p_b1 = tf.add_paragraph()
p_b1.text = "• Header Hardening: Enable CSP & X-Frame-Options."
p_b1.font.name = "Segoe UI"
p_b1.font.size = Pt(15) 
p_b1.font.color.rgb = TEXT_WHITE
p_b2 = tf.add_paragraph()
p_b2.text = "• Auto updates: Configure automatic CMS updates."
p_b2.font.name = "Segoe UI"
p_b2.font.size = Pt(15) 
p_b2.font.color.rgb = TEXT_WHITE
p_b3 = tf.add_paragraph()
p_b3.text = "• Scanning: Run monthly scans for directory leaks."
p_b3.font.name = "Segoe UI"
p_b3.font.size = Pt(15) 
p_b3.font.color.rgb = TEXT_WHITE
draw_premium_card(slide8, x_r, y_b, w_q, h_q, accent_color=YELLOW_WARN)
tb_q4 = slide8.shapes.add_textbox(Inches(x_r + 0.2), Inches(y_b + 0.2), Inches(w_q - 0.4), Inches(h_q - 0.4))
tf = tb_q4.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "4. Lead Generation & Client Acquisition"
p.font.name = "Segoe UI"
p.font.size = Pt(20) 
p.font.bold = True
p.font.color.rgb = CYAN_ACCENT
p.space_after = Pt(8)
p_b1 = tf.add_paragraph()
p_b1.text = "• Lead Gen: Pitch audit reports to low-scoring sites."
p_b1.font.name = "Segoe UI"
p_b1.font.size = Pt(15) 
p_b1.font.color.rgb = TEXT_WHITE
p_b2 = tf.add_paragraph()
p_b2.text = "• Quick Wins: Offer free SSL & GBP verification setups."
p_b2.font.name = "Segoe UI"
p_b2.font.size = Pt(15) 
p_b2.font.color.rgb = TEXT_WHITE
p_b3 = tf.add_paragraph()
p_b3.text = "• Bundles: Package security audits for CA/clinics."
p_b3.font.name = "Segoe UI"
p_b3.font.size = Pt(15) 
p_b3.font.color.rgb = TEXT_WHITE
slide9 = prs.slides.add_slide(slide_layout)
apply_background(slide9)
create_title(slide9, "Lead Conversion & Outreach Playbook", "Commercial strategy for IT consulting & cybersecurity agencies")
col_w = 4.4
gap = 0.4
y_pos = 2.4
h_pos = 5.6
draw_premium_card(slide9, 1.0, y_pos, col_w, h_pos, accent_color=RED_RISK)
tb_c1 = slide9.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.3), Inches(col_w - 0.4), Inches(h_pos - 0.6))
tf = tb_c1.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "01  |  Lead Prioritization"
p.font.name = "Segoe UI"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = RED_RISK
p.space_after = Pt(16)
p_body1_hdr = tf.add_paragraph()
p_body1_hdr.text = "Vulnerability Value"
p_body1_hdr.font.name = "Segoe UI"
p_body1_hdr.font.size = Pt(20)
p_body1_hdr.font.bold = True
p_body1_hdr.font.color.rgb = TEXT_WHITE
p_body1 = tf.add_paragraph()
p_body1.text = "• Target critical risks (score 0-2).\n• Focus on high-review businesses.\n• High reputation = high leverage."
p_body1.font.name = "Segoe UI"
p_body1.font.size = Pt(18)
p_body1.font.color.rgb = TEXT_MUTED
p_body1.space_before = Pt(4)
p_body1.space_after = Pt(16)
p_body2_hdr = tf.add_paragraph()
p_body2_hdr.text = "Sabotage Risks"
p_body2_hdr.font.name = "Segoe UI"
p_body2_hdr.font.size = Pt(20)
p_body2_hdr.font.bold = True
p_body2_hdr.font.color.rgb = TEXT_WHITE
p_body2 = tf.add_paragraph()
p_body2.text = "• Highlight unclaimed Maps listings.\n• Stress unauthorized edit risk."
p_body2.font.name = "Segoe UI"
p_body2.font.size = Pt(18)
p_body2.font.color.rgb = TEXT_MUTED
p_body2.space_before = Pt(4)
draw_premium_card(slide9, 1.0 + col_w + gap, y_pos, col_w, h_pos, accent_color=CYAN_ACCENT)
tb_c2 = slide9.shapes.add_textbox(Inches(1.2 + col_w + gap), Inches(y_pos + 0.3), Inches(col_w - 0.4), Inches(h_pos - 0.6))
tf = tb_c2.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "02  |  Conversion Hooks"
p.font.name = "Segoe UI"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = CYAN_ACCENT
p.space_after = Pt(16)
p_body1_hdr = tf.add_paragraph()
p_body1_hdr.text = "Active Visual Warning"
p_body1_hdr.font.name = "Segoe UI"
p_body1_hdr.font.size = Pt(20)
p_body1_hdr.font.bold = True
p_body1_hdr.font.color.rgb = TEXT_WHITE
p_body1 = tf.add_paragraph()
p_body1.text = "• Show 'Not Secure' browser flags.\n• Demo exposed login directories.\n• Highlight SSL expiry warnings."
p_body1.font.name = "Segoe UI"
p_body1.font.size = Pt(18)
p_body1.font.color.rgb = TEXT_MUTED
p_body1.space_before = Pt(4)
p_body1.space_after = Pt(16)
p_body2_hdr = tf.add_paragraph()
p_body2_hdr.text = "Trust Builder Offer"
p_body2_hdr.font.name = "Segoe UI"
p_body2_hdr.font.size = Pt(20)
p_body2_hdr.font.bold = True
p_body2_hdr.font.color.rgb = TEXT_WHITE
p_body2 = tf.add_paragraph()
p_body2.text = "• Offer free Google Profile claim.\n• Charge low entry-fee for SSL setup."
p_body2.font.name = "Segoe UI"
p_body2.font.size = Pt(18)
p_body2.font.color.rgb = TEXT_MUTED
p_body2.space_before = Pt(4)
draw_premium_card(slide9, 1.0 + 2*(col_w + gap), y_pos, col_w, h_pos, accent_color=GREEN_SECURE)
tb_c3 = slide9.shapes.add_textbox(Inches(1.2 + 2*(col_w + gap)), Inches(y_pos + 0.3), Inches(col_w - 0.4), Inches(h_pos - 0.6))
tf = tb_c3.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "03  |  Upsell Strategy"
p.font.name = "Segoe UI"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = GREEN_SECURE
p.space_after = Pt(16)
p_body1_hdr = tf.add_paragraph()
p_body1_hdr.text = "Compliance Focus"
p_body1_hdr.font.name = "Segoe UI"
p_body1_hdr.font.size = Pt(20)
p_body1_hdr.font.bold = True
p_body1_hdr.font.color.rgb = TEXT_WHITE
p_body1 = tf.add_paragraph()
p_body1.text = "• Pitch compliance audits to clinics.\n• Pitch secure bookings to hotels."
p_body1.font.name = "Segoe UI"
p_body1.font.size = Pt(18)
p_body1.font.color.rgb = TEXT_MUTED
p_body1.space_before = Pt(4)
p_body1.space_after = Pt(16)
p_body2_hdr = tf.add_paragraph()
p_body2_hdr.text = "Retainer Packages"
p_body2_hdr.font.name = "Segoe UI"
p_body2_hdr.font.size = Pt(20)
p_body2_hdr.font.bold = True
p_body2_hdr.font.color.rgb = TEXT_WHITE
p_body2 = tf.add_paragraph()
p_body2.text = "• Bundle weekly scanner audits.\n• Offer monthly backups & updates."
p_body2.font.name = "Segoe UI"
p_body2.font.size = Pt(18)
p_body2.font.color.rgb = TEXT_MUTED
p_body2.space_before = Pt(4)
prs.save("National_Digital_Risk_Audit_Presentation.pptx")
print("PowerPoint presentation National_Digital_Risk_Audit_Presentation.pptx generated successfully!")

# ----------------- HTML GENERATION -----------------
import json
html_template = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Pan-India Digital Security Audit - Interactive Dashboard</title>
    <!-- Google Fonts -->
    <link href="https://fonts.googleapis.com/css2?family=Outfit:wght@300;400;600;800&family=Inter:wght@400;500;600;700&display=swap" rel="stylesheet">
    <style>
        :root {
            --bg-main: #120f1d;
            --bg-card: #1b1829;
            --bg-card-hover: #231f33;
            --border-muted: #312c45;
            --text-primary: #faf6f0;
            --text-secondary: #d5c3c6;
            --text-muted: #8b7d8b;
            --accent-cyan: #e5a88b;
            --accent-cyan-glow: rgba(229, 168, 139, 0.25);
            --color-secure: #88d49e;
            --color-secure-glow: rgba(136, 212, 158, 0.15);
            --color-warning: #f1c40f;
            --color-danger: #f07167;
            --color-danger-glow: rgba(240, 113, 103, 0.15);
            
            /* Risk tier specific colors */
            --c-very-secure: #88d49e;
            --c-secure: #2e8b57;
            --c-moderate: #f1c40f;
            --c-high: #e67e22;
            --c-very-high: #f07167;
        }

        * {
            box-sizing: border-box;
            margin: 0;
            padding: 0;
        }

        body {
            background-color: var(--bg-main);
            color: var(--text-primary);
            font-family: 'Inter', sans-serif;
            min-height: 100vh;
            overflow-x: hidden;
            background-image: 
                radial-gradient(circle at 10% 20%, rgba(229, 168, 139, 0.05) 0%, transparent 40%),
                radial-gradient(circle at 90% 80%, rgba(213, 195, 198, 0.02) 0%, transparent 40%);
        }

        h1, h2, h3, .font-display {
            font-family: 'Outfit', sans-serif;
        }

        /* Scrollbar styling */
        ::-webkit-scrollbar {
            width: 8px;
            height: 8px;
        }
        ::-webkit-scrollbar-track {
            background: var(--bg-main);
        }
        ::-webkit-scrollbar-thumb {
            background: var(--border-muted);
            border-radius: 4px;
        }
        ::-webkit-scrollbar-thumb:hover {
            background: var(--accent-cyan);
        }

        /* Container Layout */
        .wrapper {
            max-width: 1560px;
            margin: 0 auto;
            padding: 24px;
        }

        /* Header section */
        header {
            display: flex;
            justify-content: space-between;
            align-items: center;
            border-bottom: 1px solid var(--border-muted);
            padding-bottom: 20px;
            margin-bottom: 30px;
        }

        .header-title h1 {
            font-size: 32px;
            font-weight: 800;
            letter-spacing: -0.5px;
            background: linear-gradient(90deg, #ffffff 0%, var(--accent-cyan) 100%);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
        }

        .header-title p {
            color: var(--text-secondary);
            font-size: 14px;
            margin-top: 4px;
        }

        .header-badge {
            background: rgba(229, 168, 139, 0.08);
            border: 1px solid var(--accent-cyan);
            border-radius: 20px;
            padding: 6px 16px;
            font-size: 11px;
            font-weight: 700;
            color: var(--accent-cyan);
            letter-spacing: 1px;
            text-transform: uppercase;
            box-shadow: 0 0 10px rgba(229, 168, 139, 0.15);
        }

        /* Main Workspace Grid */
        .top-row-grid {
            display: grid;
            grid-template-columns: 2fr 1fr;
            gap: 20px;
            margin-bottom: 30px;
        }

        @media (max-width: 1024px) {
            .top-row-grid {
                grid-template-columns: 1fr;
            }
        }

        /* KPI Grid */
        .kpi-grid {
            display: grid;
            grid-template-columns: repeat(2, 1fr);
            gap: 16px;
            height: 100%;
        }

        @media (max-width: 480px) {
            .kpi-grid {
                grid-template-columns: 1fr;
            }
        }

        .kpi-card {
            background: var(--bg-card);
            border: 1px solid var(--border-muted);
            border-radius: 12px;
            padding: 24px;
            position: relative;
            overflow: hidden;
            display: flex;
            flex-direction: column;
            justify-content: center;
            transition: all 0.3s ease;
        }

        .kpi-card::before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            width: 100%;
            height: 4px;
            background: var(--accent-cyan);
        }

        .kpi-card.secure::before { background: var(--color-secure); }
        .kpi-card.critical::before { background: var(--color-danger); }

        .kpi-card:hover {
            transform: translateY(-2px);
            box-shadow: 0 6px 20px rgba(0, 0, 0, 0.4);
            border-color: rgba(229, 168, 139, 0.2);
        }

        .kpi-value {
            font-size: 40px;
            font-weight: 800;
            margin-top: 8px;
            font-family: 'Outfit', sans-serif;
            line-height: 1;
        }

        .kpi-label {
            color: var(--text-secondary);
            font-size: 11px;
            font-weight: 700;
            letter-spacing: 0.8px;
            text-transform: uppercase;
        }

        /* Chart Panel Card */
        .chart-panel-card {
            background: var(--bg-card);
            border: 1px solid var(--border-muted);
            border-radius: 12px;
            padding: 24px;
            display: flex;
            flex-direction: column;
            justify-content: center;
            align-items: center;
            min-height: 280px;
            position: relative;
        }

        .chart-title {
            font-size: 13px;
            font-weight: 700;
            text-transform: uppercase;
            letter-spacing: 0.8px;
            color: var(--accent-cyan);
            margin-bottom: 16px;
            width: 100%;
            text-align: center;
        }

        /* SVG Donut Chart styling */
        .donut-container {
            position: relative;
            width: 160px;
            height: 160px;
        }

        .donut-svg {
            transform: rotate(-90deg);
        }

        .donut-ring {
            fill: none;
            stroke: var(--border-muted);
            stroke-width: 10;
        }

        .donut-segment {
            fill: none;
            stroke-width: 12;
            transition: stroke-dasharray 0.4s ease, stroke-dashoffset 0.4s ease;
        }

        .donut-center-text {
            position: absolute;
            top: 50%;
            left: 50%;
            transform: translate(-50%, -50%);
            text-align: center;
            display: flex;
            flex-direction: column;
            justify-content: center;
            align-items: center;
        }

        .donut-center-num {
            font-family: 'Outfit', sans-serif;
            font-size: 26px;
            font-weight: 800;
            color: var(--accent-cyan);
            line-height: 1;
        }

        .donut-center-lbl {
            font-size: 9px;
            color: var(--text-secondary);
            font-weight: 700;
            text-transform: uppercase;
            margin-top: 2px;
        }

        /* Chart Legends */
        .chart-legend {
            display: flex;
            flex-wrap: wrap;
            justify-content: center;
            gap: 12px;
            margin-top: 16px;
            font-size: 11px;
            color: var(--text-secondary);
        }

        .legend-item {
            display: flex;
            align-items: center;
            gap: 6px;
        }

        .legend-dot {
            width: 8px;
            height: 8px;
            border-radius: 50%;
        }

        /* Controls Section (Filters & Search) */
        .controls-panel {
            background: var(--bg-card);
            border: 1px solid var(--border-muted);
            border-radius: 12px;
            padding: 20px;
            margin-bottom: 30px;
            display: flex;
            flex-direction: column;
            gap: 16px;
        }

        .search-row {
            display: flex;
            gap: 16px;
            flex-wrap: wrap;
        }

        .search-box {
            flex: 1;
            min-width: 280px;
            position: relative;
        }

        .search-box input {
            width: 100%;
            background: var(--bg-main);
            border: 1px solid var(--border-muted);
            border-radius: 8px;
            padding: 12px 16px;
            color: var(--text-primary);
            font-family: 'Inter', sans-serif;
            font-size: 14px;
            transition: all 0.2s ease;
        }

        .search-box input:focus {
            outline: none;
            border-color: var(--accent-cyan);
            box-shadow: 0 0 10px rgba(229, 168, 139, 0.15);
        }

        .filter-row {
            display: flex;
            justify-content: space-between;
            align-items: center;
            flex-wrap: wrap;
            gap: 16px;
        }

        .filter-group {
            display: flex;
            align-items: center;
            gap: 8px;
            flex-wrap: wrap;
        }

        .filter-label {
            color: var(--text-secondary);
            font-size: 12px;
            font-weight: 600;
            margin-right: 4px;
        }

        .btn {
            background: var(--bg-main);
            border: 1px solid var(--border-muted);
            border-radius: 6px;
            padding: 8px 14px;
            color: var(--text-secondary);
            font-size: 12px;
            font-weight: 600;
            cursor: pointer;
            transition: all 0.2s ease;
        }

        .btn:hover {
            color: var(--text-primary);
            border-color: var(--text-secondary);
        }

        .btn.active {
            background: rgba(229, 168, 139, 0.08);
            border-color: var(--accent-cyan);
            color: var(--accent-cyan);
        }

        /* View Toggle Mode switcher */
        .view-switcher {
            display: flex;
            background: var(--bg-main);
            border: 1px solid var(--border-muted);
            border-radius: 8px;
            padding: 4px;
        }

        .view-btn {
            background: transparent;
            border: none;
            border-radius: 6px;
            padding: 6px 12px;
            color: var(--text-secondary);
            font-size: 12px;
            font-weight: 600;
            cursor: pointer;
            transition: all 0.2s ease;
        }

        .view-btn.active {
            background: var(--border-muted);
            color: var(--text-primary);
        }

        /* Lead Cards View */
        .leads-grid {
            display: grid;
            grid-template-columns: repeat(auto-fill, minmax(360px, 1fr));
            gap: 20px;
        }

        .lead-card {
            background: var(--bg-card);
            border: 1px solid var(--border-muted);
            border-radius: 12px;
            padding: 24px;
            cursor: pointer;
            transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
            position: relative;
        }

        .lead-card::after {
            content: '';
            position: absolute;
            inset: 0;
            border-radius: 12px;
            border: 1px solid var(--accent-cyan);
            opacity: 0;
            transition: all 0.3s ease;
            pointer-events: none;
        }

        .lead-card:hover {
            transform: translateY(-4px);
            box-shadow: 0 10px 30px rgba(229, 168, 139, 0.06);
            border-color: transparent;
        }

        .lead-card:hover::after {
            opacity: 0.35;
        }

        .card-header {
            display: flex;
            justify-content: space-between;
            align-items: flex-start;
            margin-bottom: 16px;
        }

        .card-title h3 {
            font-size: 18px;
            font-weight: 700;
            color: var(--text-primary);
            line-height: 1.2;
        }

        .card-title p {
            color: var(--text-secondary);
            font-size: 11px;
            margin-top: 4px;
            text-transform: uppercase;
            font-weight: 600;
        }

        .score-badge {
            background: rgba(255, 255, 255, 0.05);
            border: 1px solid var(--border-muted);
            border-radius: 8px;
            padding: 6px 12px;
            display: flex;
            flex-direction: column;
            align-items: center;
        }

        .score-val {
            font-size: 18px;
            font-weight: 800;
            font-family: 'Outfit', sans-serif;
        }

        .score-val.secure { color: var(--color-secure); }
        .score-val.warning { color: var(--color-warning); }
        .score-val.danger { color: var(--color-danger); }

        .score-label {
            font-size: 9px;
            color: var(--text-secondary);
            margin-top: 2px;
            font-weight: 600;
        }

        .card-body {
            margin-bottom: 16px;
        }

        .meta-item {
            display: flex;
            justify-content: space-between;
            font-size: 13px;
            margin-bottom: 6px;
        }

        .meta-label {
            color: var(--text-secondary);
        }

        .meta-val {
            color: var(--text-primary);
            font-weight: 500;
        }

        .priority-indicator {
            border-radius: 4px;
            padding: 2px 8px;
            font-size: 10px;
            font-weight: 700;
            letter-spacing: 0.5px;
        }

        .priority-indicator.critical {
            background: rgba(239, 68, 68, 0.1);
            color: var(--color-danger);
            border: 1px solid rgba(239, 68, 68, 0.2);
        }
        .priority-indicator.medium {
            background: rgba(245, 158, 11, 0.1);
            color: var(--color-warning);
            border: 1px solid rgba(245, 158, 11, 0.2);
        }
        .priority-indicator.low {
            background: rgba(16, 185, 129, 0.1);
            color: var(--color-secure);
            border: 1px solid rgba(16, 185, 129, 0.2);
        }

        .card-footer {
            border-top: 1px solid var(--border-muted);
            padding-top: 12px;
            font-size: 12px;
            color: var(--accent-cyan);
            font-weight: 600;
            overflow: hidden;
            text-overflow: ellipsis;
            white-space: nowrap;
        }

        /* Database Table View */
        .table-container {
            background: var(--bg-card);
            border: 1px solid var(--border-muted);
            border-radius: 12px;
            overflow-x: auto;
            display: none; /* Switched by Javascript */
        }

        table {
            width: 100%;
            border-collapse: collapse;
            text-align: left;
            font-size: 13px;
        }

        th {
            background: #111524;
            color: var(--accent-cyan);
            font-weight: 700;
            padding: 16px;
            border-bottom: 1px solid var(--border-muted);
            font-size: 11px;
            letter-spacing: 0.8px;
            text-transform: uppercase;
        }

        td {
            padding: 14px 16px;
            border-bottom: 1px solid var(--border-muted);
            vertical-align: middle;
            color: var(--text-primary);
        }

        tr:last-child td {
            border-bottom: none;
        }

        tr:hover td {
            background: rgba(255, 255, 255, 0.02);
            cursor: pointer;
        }

        tr.zebra td {
            background: rgba(255, 255, 255, 0.008);
        }

        /* Detail Modal */
        .modal-overlay {
            position: fixed;
            inset: 0;
            background: rgba(18, 15, 29, 0.85);
            backdrop-filter: blur(12px);
            z-index: 1000;
            display: flex;
            justify-content: center;
            align-items: center;
            opacity: 0;
            pointer-events: none;
            transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
        }

        .modal-overlay.active {
            opacity: 1;
            pointer-events: all;
        }

        .modal-content {
            background: var(--bg-card);
            border: 1px solid var(--accent-cyan);
            border-radius: 16px;
            width: 90%;
            max-width: 820px;
            max-height: 85vh;
            overflow-y: auto;
            padding: 32px;
            transform: scale(0.95);
            transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
            box-shadow: 0 0 50px rgba(229, 168, 139, 0.15);
            position: relative;
        }

        .modal-overlay.active .modal-content {
            transform: scale(1);
        }

        .modal-close {
            position: absolute;
            top: 24px;
            right: 24px;
            background: transparent;
            border: none;
            color: var(--text-secondary);
            font-size: 24px;
            cursor: pointer;
            transition: all 0.2s ease;
        }

        .modal-close:hover {
            color: var(--accent-cyan);
        }

        .modal-header {
            border-bottom: 1px solid var(--border-muted);
            padding-bottom: 18px;
            margin-bottom: 24px;
        }

        .modal-header h2 {
            font-size: 28px;
            font-weight: 800;
            color: #ffffff;
        }

        .modal-category {
            color: var(--accent-cyan);
            font-size: 13px;
            font-weight: 600;
            text-transform: uppercase;
            margin-top: 4px;
        }

        .modal-section {
            margin-bottom: 24px;
        }

        .modal-section-title {
            color: var(--accent-cyan);
            font-size: 14px;
            font-weight: 700;
            letter-spacing: 0.8px;
            text-transform: uppercase;
            margin-bottom: 12px;
        }

        .vulnerability-badge-list {
            display: flex;
            flex-wrap: wrap;
            gap: 8px;
            margin-top: 10px;
        }

        .vulnerability-tag {
            background: rgba(239, 68, 68, 0.08);
            border: 1px solid rgba(239, 68, 68, 0.3);
            border-radius: 6px;
            padding: 6px 12px;
            font-size: 12px;
            color: var(--color-danger);
            font-weight: 500;
        }

        .vulnerability-tag.secure {
            background: rgba(16, 185, 129, 0.08);
            border: 1px solid rgba(16, 185, 129, 0.3);
            color: var(--color-secure);
        }

        .recs-box {
            background: var(--bg-main);
            border: 1px solid var(--border-muted);
            border-radius: 8px;
            padding: 16px 20px;
            font-size: 14px;
            line-height: 1.6;
            color: var(--text-primary);
        }

        .pitch-box {
            background: rgba(229, 168, 139, 0.05);
            border: 1px solid rgba(229, 168, 139, 0.2);
            border-radius: 8px;
            padding: 16px 20px;
            font-size: 14px;
            line-height: 1.5;
            color: var(--accent-cyan);
            font-weight: 500;
        }

        .details-grid {
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(220px, 1fr));
            gap: 16px;
        }

        .detail-item {
            background: var(--bg-main);
            border: 1px solid var(--border-muted);
            border-radius: 8px;
            padding: 12px 16px;
        }

        .detail-label {
            font-size: 11px;
            color: var(--text-secondary);
            text-transform: uppercase;
            font-weight: 600;
        }

        .detail-value {
            font-size: 15px;
            font-weight: 600;
            color: #ffffff;
            margin-top: 4px;
        }

        .detail-value.red { color: var(--color-danger); }
        .detail-value.green { color: var(--color-secure); }
        .detail-value.yellow { color: var(--color-warning); }

        /* Empty search state */
        .empty-state {
            background: var(--bg-card);
            border: 1px solid var(--border-muted);
            border-radius: 12px;
            padding: 50px 20px;
            text-align: center;
            color: var(--text-secondary);
            font-size: 16px;
            display: none;
        }

        /* Clickable website links */
        a {
            transition: opacity 0.2s ease;
        }
        a:hover {
            opacity: 0.85;
            text-decoration: underline !important;
        }

        /* Premium layout updates for 3-column top row grid */
        .top-row-grid {
            grid-template-columns: 1.2fr 1fr 1fr !important;
        }

        @media (max-width: 1200px) {
            .top-row-grid {
                grid-template-columns: 1fr !important;
            }
        }

        /* SVG Bar Chart styling */
        .bar-chart-container {
            width: 100%;
            height: 160px;
            display: flex;
            flex-direction: column;
            justify-content: space-between;
        }

        .bar-row {
            display: flex;
            align-items: center;
            font-size: 11px;
            gap: 8px;
        }

        .bar-label {
            width: 110px;
            color: var(--text-secondary);
            font-weight: 600;
            white-space: nowrap;
            overflow: hidden;
            text-overflow: ellipsis;
        }

        .bar-track {
            flex: 1;
            height: 12px;
            background: var(--border-muted);
            border-radius: 6px;
            overflow: hidden;
            position: relative;
        }

        .bar-fill {
            height: 100%;
            border-radius: 6px;
            transition: width 0.4s ease;
            background: linear-gradient(90deg, var(--accent-cyan) 0%, var(--c-very-secure) 100%);
        }

        .bar-val {
            width: 32px;
            text-align: right;
            font-weight: 700;
            color: var(--accent-cyan);
        }

        /* Executive Briefing Report Section Styles */
        .report-container {
            max-width: 1000px;
            margin: 0 auto;
            display: flex;
            flex-direction: column;
            gap: 30px;
            padding: 20px 0;
        }

        .report-section {
            background: var(--bg-card);
            border: 1px solid var(--border-muted);
            border-radius: 16px;
            padding: 32px;
            box-shadow: 0 4px 30px rgba(0, 0, 0, 0.2);
            transition: all 0.3s ease;
        }

        .report-section:hover {
            border-color: rgba(229, 168, 139, 0.15);
        }

        .report-title-row {
            display: flex;
            align-items: center;
            gap: 12px;
            margin-bottom: 20px;
            border-bottom: 1px solid var(--border-muted);
            padding-bottom: 12px;
        }

        .report-icon {
            font-size: 24px;
            color: var(--accent-cyan);
        }

        .report-h2 {
            font-family: 'Outfit', sans-serif;
            font-size: 22px;
            font-weight: 700;
            color: #ffffff;
            letter-spacing: -0.5px;
        }

        .report-text {
            color: var(--text-secondary);
            font-size: 14.5px;
            line-height: 1.7;
            margin-bottom: 16px;
        }

        .report-bullet {
            margin-left: 20px;
            margin-bottom: 10px;
            color: var(--text-secondary);
            font-size: 14px;
            line-height: 1.6;
        }

        /* Alert block in report */
        .report-alert {
            background: rgba(240, 113, 103, 0.05);
            border-left: 4px solid var(--color-danger);
            border-radius: 0 8px 8px 0;
            padding: 16px 20px;
            margin: 20px 0;
            color: var(--text-primary);
        }

        .report-alert strong {
            color: var(--color-danger);
            font-family: 'Outfit', sans-serif;
            font-size: 15px;
            display: block;
            margin-bottom: 6px;
        }

        /* Dynamic stat grid inside report */
        .report-stat-grid {
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
            gap: 16px;
            margin: 20px 0;
        }

        .report-stat-card {
            background: var(--bg-main);
            border: 1px solid var(--border-muted);
            border-radius: 10px;
            padding: 20px;
            text-align: center;
        }

        .report-stat-val {
            font-size: 24px;
            font-weight: 800;
            color: var(--accent-cyan);
            font-family: 'Outfit', sans-serif;
        }

        .report-stat-lbl {
            font-size: 11px;
            text-transform: uppercase;
            letter-spacing: 0.8px;
            color: var(--text-muted);
            margin-top: 4px;
            font-weight: 600;
        }


        /* City pills styling */
        .city-pills-row {
            display: flex;
            gap: 8px;
            overflow-x: auto;
            padding: 4px 0 12px 0;
            margin-bottom: 20px;
            border-bottom: 1px solid var(--border-muted);
            scrollbar-width: thin;
        }
        
        .city-pill {
            background: var(--bg-card);
            border: 1px solid var(--border-muted);
            color: var(--text-secondary);
            padding: 6px 14px;
            border-radius: 20px;
            font-size: 12px;
            font-weight: 600;
            cursor: pointer;
            white-space: nowrap;
            transition: all 0.2s ease;
        }
        
        .city-pill:hover {
            border-color: var(--accent-cyan);
            color: #ffffff;
            transform: translateY(-1px);
        }
        
        .city-pill.active {
            background: linear-gradient(135deg, var(--accent-cyan) 0%, var(--c-high) 100%);
            border-color: transparent;
            color: #120f1d;
            box-shadow: 0 4px 12px rgba(229, 168, 139, 0.2);
        }

        /* Floating drawer styling */
        .drawer-overlay {
            position: fixed;
            top: 0;
            left: 0;
            width: 100%;
            height: 100%;
            background: rgba(0, 0, 0, 0.6);
            backdrop-filter: blur(4px);
            z-index: 1999;
            display: none;
            opacity: 0;
            transition: opacity 0.3s ease;
        }

        .calc-drawer {
            position: fixed;
            top: 0;
            right: -420px;
            width: 400px;
            height: 100%;
            background: var(--bg-card);
            border-left: 1px solid var(--border-muted);
            box-shadow: -10px 0 30px rgba(0, 0, 0, 0.6);
            z-index: 2000;
            transition: right 0.3s cubic-bezier(0.16, 1, 0.3, 1);
            padding: 30px;
            display: flex;
            flex-direction: column;
            gap: 20px;
            overflow-y: auto;
        }

        .calc-drawer.open {
            right: 0;
        }

        .calc-drawer-header {
            display: flex;
            justify-content: space-between;
            align-items: center;
            border-bottom: 1px solid var(--border-muted);
            padding-bottom: 16px;
        }

        .calc-title {
            font-family: 'Outfit', sans-serif;
            font-size: 20px;
            font-weight: 700;
            color: #ffffff;
        }

        .calc-close-btn {
            background: none;
            border: none;
            color: var(--text-muted);
            font-size: 20px;
            cursor: pointer;
            transition: color 0.2s;
        }

        .calc-close-btn:hover {
            color: #ffffff;
        }

        .calc-input-group {
            display: flex;
            flex-direction: column;
            gap: 8px;
        }

        .calc-label {
            font-size: 12px;
            color: var(--text-secondary);
            font-weight: 600;
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }

        .calc-input {
            background: var(--bg-main);
            border: 1px solid var(--border-muted);
            color: #ffffff;
            padding: 10px 14px;
            border-radius: 8px;
            font-size: 14px;
            font-family: 'Inter', sans-serif;
        }

        .calc-input:focus {
            outline: none;
            border-color: var(--accent-cyan);
        }

        .checkbox-label-row {
            display: flex;
            align-items: flex-start;
            gap: 12px;
            padding: 12px;
            background: var(--bg-main);
            border: 1px solid var(--border-muted);
            border-radius: 8px;
            cursor: pointer;
            transition: all 0.2s ease;
            user-select: none;
        }

        .checkbox-label-row:hover {
            border-color: rgba(229, 168, 139, 0.15);
            background: rgba(255, 255, 255, 0.01);
        }

        .checkbox-label-row input[type="checkbox"] {
            margin-top: 3px;
            accent-color: var(--accent-cyan);
        }

        .checkbox-text {
            display: flex;
            flex-direction: column;
            gap: 2px;
        }

        .checkbox-title {
            font-size: 13.5px;
            font-weight: 600;
            color: #ffffff;
        }

        .checkbox-desc {
            font-size: 11px;
            color: var(--text-muted);
        }

        .calc-result-box {
            background: rgba(229, 168, 139, 0.04);
            border: 1px solid rgba(229, 168, 139, 0.15);
            border-radius: 12px;
            padding: 20px;
            display: flex;
            flex-direction: column;
            align-items: center;
            gap: 14px;
            margin-top: 10px;
        }

        .calc-meter-container {
            position: relative;
            width: 100px;
            height: 100px;
            display: flex;
            justify-content: center;
            align-items: center;
        }

        .calc-score-display {
            font-family: 'Outfit', sans-serif;
            font-size: 32px;
            font-weight: 800;
            color: var(--accent-cyan);
        }

        .calc-result-metrics {
            width: 100%;
            display: flex;
            flex-direction: column;
            gap: 6px;
            font-size: 12px;
            border-top: 1px solid var(--border-muted);
            padding-top: 10px;
        }

        .calc-metric-row {
            display: flex;
            justify-content: space-between;
        }

        .calc-metric-lbl {
            color: var(--text-muted);
        }

        .calc-metric-val {
            font-weight: 700;
            color: #ffffff;
        }

    
        @media print {
            body { background: #ffffff !important; color: #000000 !important; }
            .header-title h1, .report-h2, .report-stat-val { color: #000000 !important; }
            header { background: #ffffff !important; border-bottom: 2px solid #000; padding: 10px; }
            .controls-panel, .view-switcher, #leads-section-container, #dashboard-top-grid, .btn { display: none !important; }
            #report-section-container { display: block !important; }
            .report-section { box-shadow: none; border: 1px solid #ccc; break-inside: avoid; margin-bottom: 20px; page-break-inside: avoid; }
        }

        body.light-mode {
            --bg-dark: #f0f2f5;
            --bg-card: #ffffff;
            --text-light: #1e293b;
            --text-muted: #64748b;
            --border-muted: #cbd5e1;
            --header-bg: #ffffff;
        }
        body.light-mode header {
            box-shadow: 0 4px 15px rgba(0,0,0,0.05);
            border-bottom: none;
        }
        body.light-mode .lead-card, body.light-mode .report-section {
            box-shadow: 0 4px 12px rgba(0,0,0,0.05);
            border: 1px solid var(--border-muted);
        }
        body.light-mode .stat-card {
            background: #ffffff;
            box-shadow: 0 4px 12px rgba(0,0,0,0.05);
            border: 1px solid var(--border-muted);
        }
        body.light-mode th {
            background-color: #f8fafc;
            color: #475569;
            border-bottom-color: #cbd5e1;
        }
        body.light-mode td {
            border-bottom-color: #e2e8f0;
        }
        body.light-mode tr:hover {
            background-color: #f1f5f9;
        }
        body.light-mode .calc-drawer {
            background: #ffffff;
            box-shadow: -5px 0 25px rgba(0,0,0,0.1);
        }
        body.light-mode .calc-drawer-header {
            border-bottom-color: #e2e8f0;
        }
        body.light-mode .donut-text-val {
            fill: #1e293b;
        }
        body.light-mode .donut-text-label {
            fill: #64748b;
        }
        body.light-mode .modal-content {
            background: #ffffff;
            color: #1e293b;
        }
        body.light-mode .modal-header {
            border-bottom-color: #e2e8f0;
        }
        body.light-mode .email-pitch-box {
            background: #f8fafc;
            border-color: #cbd5e1;
            color: #334155;
        }

    </style>
</head>
<body>

    <div class="wrapper">
        <header>
            <div class="header-title">
                <h1>National SMB Digital Risk Auditor</h1>
                <p>Digital security posture & vulnerability assessment of 50 Indian commercial leads</p>
            </div>
            <div style="display: flex; align-items: center; gap: 16px;">
                <div class="view-switcher" style="padding: 2px;">
                    <button class="view-btn active" id="tab-leads-btn">Audit Database</button>
                    <button class="view-btn" id="tab-report-btn">Executive Briefing</button>
                </div>
                <button class="btn" id="theme-toggle-btn" style="background: rgba(229, 168, 139, 0.1); border-color: var(--color-warning); color: var(--color-warning);">☀️ Light Mode</button>
                <button class="btn" id="open-calc-btn" style="background: rgba(136, 212, 158, 0.08); border-color: var(--color-secure); color: var(--color-secure); display: flex; align-items: center; gap: 6px;">
                    🧮 Posture Calculator
                </button>
                <div class="header-badge">PAN-INDIA AUDIT • 2026</div>
            </div>
        </header>

        <div class="top-row-grid" id="dashboard-top-grid">
            <!-- KPI Summary Cards -->
            <div class="kpi-grid">
                <div class="kpi-card">
                    <div class="kpi-label">Businesses Checked</div>
                    <div class="kpi-value" id="kpi-total">50</div>
                </div>
                <div class="kpi-card">
                    <div class="kpi-label">Average Security Score</div>
                    <div class="kpi-value" id="kpi-score" style="color: var(--accent-cyan);">3.44</div>
                </div>
                <div class="kpi-card secure">
                    <div class="kpi-label">Secure Rate (Score 4-5)</div>
                    <div class="kpi-value" id="kpi-secure" style="color: var(--color-secure);">64.0%</div>
                </div>
                <div class="kpi-card critical">
                    <div class="kpi-label">Critical At-Risk Leads</div>
                    <div class="kpi-value" id="kpi-risky" style="color: var(--color-danger);">18</div>
                </div>
            </div>
            
            <!-- Live Donut Chart Panel -->
            <div class="chart-panel-card">
                <div class="chart-title">Risk Level Breakdown (Live)</div>
                <div class="donut-container">
                    <svg class="donut-svg" width="160" height="160" viewBox="0 0 160 160">
                        <circle class="donut-ring" cx="80" cy="80" r="60"></circle>
                        <!-- Donut slices generated dynamically -->
                        <g id="donut-segments-group"></g>
                    </svg>
                    <div class="donut-center-text">
                        <span class="donut-center-num" id="donut-center-count">50</span>
                        <span class="donut-center-lbl">Leads</span>
                    </div>
                </div>
                <div class="chart-legend">
                    <div class="legend-item"><span class="legend-dot" style="background: var(--c-very-secure);"></span> Very Secure</div>
                    <div class="legend-item"><span class="legend-dot" style="background: var(--c-secure);"></span> Secure</div>
                    <div class="legend-item"><span class="legend-dot" style="background: var(--c-high);"></span> High Risk</div>
                    <div class="legend-item"><span class="legend-dot" style="background: var(--c-very-high);"></span> Very High</div>
                </div>
            </div>

            <!-- Live Category Scores Bar Chart Panel -->
            <div class="chart-panel-card">
                <div class="chart-title">Score by Category (Live)</div>
                <div class="bar-chart-container" id="bar-chart-container">
                    <!-- Populated dynamically via Javascript -->
                </div>
            </div>
        </div>

        <div id="leads-section-container">
        <!-- City Filter Pills -->
        <div class="city-pills-row" id="city-pills-row">
            <!-- Populated dynamically via Javascript -->
        </div>
        <!-- Filters Section -->
        <div class="controls-panel">
            <div class="search-row">
                <div class="search-box">
                    
            <select class="filter-select" id="sort-filter">
                <option value="score-asc">Sort: Risk (High to Low)</option>
                <option value="score-desc">Sort: Score (High to Low)</option>
                <option value="name-asc">Sort: Name (A-Z)</option>
                <option value="reviews-desc">Sort: Popularity (Reviews)</option>
            </select>

            <input type="text" id="search-input" placeholder="Search business name, location, or vulnerability...">
                </div>
            </div>
            <div class="filter-row">
                <div class="filter-group" id="filter-category">
                    <span class="filter-label">Category:</span>
                    <button class="btn active" data-cat="all">All</button>
                    <button class="btn" data-cat="Restaurants & Cafes">Restaurants</button>
                    <button class="btn" data-cat="Healthcare & Clinics">Healthcare</button>
                    <button class="btn" data-cat="Retail & Shops">Retail</button>
                    <button class="btn" data-cat="Professional Services">Professional Services</button>
                    <button class="btn" data-cat="Hospitality & Hotels">Hospitality</button>
                </div>
                <div class="filter-group" id="filter-priority">
                    <span class="filter-label">Priority:</span>
                    <button class="btn active" data-priority="all">All</button>
                    <button class="btn" data-priority="CRITICAL">Critical</button>
                    <button class="btn" data-priority="MEDIUM">Medium</button>
                    <button class="btn" data-priority="LOW">Low</button>
                </div>
                <div style="display: flex; gap: 12px; align-items: center;">
                    <button class="btn" id="export-csv-btn" style="background: rgba(229, 168, 139, 0.08); border-color: var(--accent-cyan); color: var(--accent-cyan); display: flex; align-items: center; gap: 6px;">
                        📥 Export Leads CSV
                    </button>
                    <div class="view-switcher">
                        <button class="view-btn active" id="view-cards-btn">Cards</button>
                        <button class="view-btn" id="view-table-btn">Database Table</button>
                    </div>
                </div>
            </div>
        </div>

        <!-- Empty Search State -->
        <div class="empty-state" id="empty-state">
            No businesses matches the selected filter criteria. Try expanding your search.
        </div>

        <!-- Cards View -->
        <div class="leads-grid" id="leads-grid">
            <!-- Populated by JavaScript -->
        </div>

        <!-- Table View -->
        <div class="table-container" id="table-container">
            <table>
                <thead>
                    <tr>
                        <th>ID</th>
                        <th>Business Name</th>
                        <th>Website Link</th>
                        <th>Category</th>
                        <th>Location</th>
                        <th>Score</th>
                        <th>Risk Level</th>
                        <th>Priority</th>
                        <th>HTTPS</th>
                        <th>SSL Status</th>
                    </tr>
                </thead>
                <tbody id="table-body">
                    <!-- Populated by JavaScript -->
                </tbody>
            </table>
        </div>
    </div>

        </div> <!-- End of leads-section-container -->

        <!-- Executive Briefing Report Section -->
        <div id="report-section-container" style="display: none;">
            <div class="report-container">
                <!-- Section 1: Executive Summary -->
                <div class="report-section">
                    <div class="report-title-row">
                        <span class="report-icon">🛡</span>
                        <h2 class="report-h2">1. Executive Summary & Assessment Overview</h2>
<button class="btn" onclick="window.print()" style="margin-left: 15px; background: rgba(229, 168, 139, 0.1); border-color: var(--accent-cyan); color: var(--accent-cyan);">🖨 Print Report</button>
                    </div>
                    <p class="report-text">
                        This auditing platform captures the digital security posture and cybersecurity exposure of <strong>50 commercial businesses</strong> distributed across primary Indian economic hubs. Sourced dynamically via Google Maps profiles, target leads span 5 main categories: <em>Restaurants & Cafes</em>, <em>Healthcare & Clinics</em>, <em>Retail & Shops</em>, <em>Professional Services</em>, and <em>Hospitality & Hotels</em>.
                    </p>
                    <div class="report-alert">
                        <strong>⚠ DIGITAL POSTURE ANALYSIS THREAT WARNING</strong>
                        A significant portion of the Indian small-to-medium business segment (36.0% of assessed leads) operates under vulnerable or unencrypted configurations (Security Score 0-3). Omission of primary transport layers and administrative protections invites risk of competitor profile takeover, brand disruption, and customer booking data hijacking.
                    </div>
                    <div class="report-stat-grid">
                        <div class="report-stat-card">
                            <div class="report-stat-val">3.44 / 5.0</div>
                            <div class="report-stat-lbl">Average security score</div>
                        </div>
                        <div class="report-stat-card">
                            <div class="report-stat-val">64.0%</div>
                            <div class="report-stat-lbl">Secure rate (score 4-5)</div>
                        </div>
                        <div class="report-stat-card">
                            <div class="report-stat-val">18 Leads</div>
                            <div class="report-stat-lbl">Critical At-Risk Leads</div>
                        </div>
                    </div>
                </div>

                <!-- Section 2: Key Findings -->
                <div class="report-section">
                    <div class="report-title-row">
                        <span class="report-icon">📊</span>
                        <h2 class="report-h2">2. Key Findings & Digital Risk Indicators</h2>
                    </div>
                    <p class="report-text">
                        Cross-sector testing across target domains revealed 56 distinct system-level security gaps:
                    </p>
                    <div class="report-bullet"><strong>• Digital Absence (16.0%):</strong> 8 businesses operate without active landing pages or domains, cutting off digital customer acquisition channels.</div>
                    <div class="report-bullet"><strong>• Unencrypted Protocols (20.0%):</strong> 10 businesses fail to force HTTPS redirects or use expired SSL configurations, transmitting booking inputs in cleartext.</div>
                    <div class="report-bullet"><strong>• Public Admin Endpoints (12.0%):</strong> 6 active websites expose administrative URL interfaces (such as <code>/wp-admin</code>), enabling brute-force automated login scans.</div>
                    <div class="report-bullet"><strong>• Unclaimed Listings (18.0%):</strong> 9 businesses leave Google Business Profiles unclaimed, allowing malicious public phone or name adjustments.</div>
                </div>

                <!-- Section 3: Industry Postures -->
                <div class="report-section">
                    <div class="report-title-row">
                        <span class="report-icon">🏢</span>
                        <h2 class="report-h2">3. Sector & Regional Performance Scores</h2>
                    </div>
                    <p class="report-text">
                        Auditing reveals wide disparity in cybersecurity implementation:
                    </p>
                    <div class="report-bullet"><strong>• Hospitality & Hotels (Avg: 4.20/5):</strong> Enterprise compliance protocols ensure HTTPS and SSL standardization.</div>
                    <div class="report-bullet"><strong>• Healthcare & Clinics (Avg: 3.70/5):</strong> Moderate score, though small dental clinics exhibit missing header vulnerabilities.</div>
                    <div class="report-bullet"><strong>• Professional Services (Avg: 3.50/5):</strong> Higher base score but minor CA offices lack form protections.</div>
                    <div class="report-bullet"><strong>• Retail & Shops (Avg: 3.30/5):</strong> Exposed local boutique sites suffer from missing encryption redirects.</div>
                    <div class="report-bullet"><strong>• Restaurants & Cafes (Avg: 2.50/5):</strong> Lowest scoring industry due to legacy local web configurations.</div>
                </div>

                <!-- Section 4: Action Roadmap -->
                <div class="report-section">
                    <div class="report-title-row">
                        <span class="report-icon">🚀</span>
                        <h2 class="report-h2">4. Actionable Remediation Roadmap</h2>
                    </div>
                    <p class="report-text">
                        Remediation should proceed in three tactical phases to secure brand presence:
                    </p>
                    <div class="report-bullet"><strong>1. Immediate:</strong> Enforce 301 HTTPS redirects, install valid SSL certificates, and verify maps ownership tags.</div>
                    <div class="report-bullet"><strong>2. Medium-Term:</strong> Restrict admin directory access by IP range or rename login endpoints; deploy anti-CSRF protection on input forms.</div>
                    <div class="report-bullet"><strong>3. Long-Term:</strong> Add HTTP Strict Transport Security (HSTS) headers, restrict directory folder listings, and schedule monthly malware scans.</div>
                </div>
            </div>
        </div>

    <!-- Score Calculator Drawer overlay & card -->
    <div class="drawer-overlay" id="drawer-overlay"></div>
    <div class="calc-drawer" id="calc-drawer">
        <div class="calc-drawer-header">
            <h3 class="calc-title">🛡 Posture Score Calculator</h3>
            <button class="calc-close-btn" id="close-calc-btn">×</button>
        </div>
        <p style="color: var(--text-secondary); font-size: 12.5px; line-height: 1.5;">
            Simulate a security assessment audit in real-time to calculate potential scores and pitches for new prospects.
        </p>

        <div class="calc-input-group">
            <label class="calc-label">Business Name</label>
            <input type="text" class="calc-input" id="calc-biz-name" placeholder="e.g. Apollo Pharmacy">
        </div>

        <div class="calc-input-group">
            <label class="calc-label">Vulnerabilities Identified</label>
            <div style="display: flex; flex-direction: column; gap: 8px; margin-top: 4px;">
                <label class="checkbox-label-row">
                    <input type="checkbox" id="calc-v-http" checked>
                    <div class="checkbox-text">
                        <span class="checkbox-title">No HTTPS Redirect</span>
                        <span class="checkbox-desc">Deducts 1.0 point. Browser shows unencrypted warning.</span>
                    </div>
                </label>
                <label class="checkbox-label-row">
                    <input type="checkbox" id="calc-v-ssl" checked>
                    <div class="checkbox-text">
                        <span class="checkbox-title">Missing/Expired SSL Certificate</span>
                        <span class="checkbox-desc">Deducts 1.0 point. Connection is flagged private.</span>
                    </div>
                </label>
                <label class="checkbox-label-row">
                    <input type="checkbox" id="calc-v-admin">
                    <div class="checkbox-text">
                        <span class="checkbox-title">Exposed Admin Panel (/wp-admin)</span>
                        <span class="checkbox-desc">Deducts 1.0 point. Susceptible to brute-force scans.</span>
                    </div>
                </label>
                <label class="checkbox-label-row">
                    <input type="checkbox" id="calc-v-forms">
                    <div class="checkbox-text">
                        <span class="checkbox-title">Directory Indexing / Unprotected Forms</span>
                        <span class="checkbox-desc">Deducts 1.0 point. Exposed files or contact form spam.</span>
                    </div>
                </label>
                <label class="checkbox-label-row">
                    <input type="checkbox" id="calc-v-safe">
                    <div class="checkbox-text">
                        <span class="checkbox-title">Safe Browsing Suspicious Warning</span>
                        <span class="checkbox-desc">Deducts 1.0 point. Flagged by Google security filters.</span>
                    </div>
                </label>
            </div>
        </div>

        <div class="calc-result-box">
            <div class="calc-label">Simulated Posture Score</div>
            <div class="calc-meter-container">
                <span class="calc-score-display" id="calc-score-val">3.0</span>
            </div>
            
            <div class="calc-result-metrics">
                <div class="calc-metric-row">
                    <span class="calc-metric-lbl">Risk Classification</span>
                    <span class="calc-metric-val" id="calc-risk-val" style="color: var(--color-danger);">Moderate Risk</span>
                </div>
                <div class="calc-metric-row">
                    <span class="calc-metric-lbl">Outreach Priority</span>
                    <span class="calc-metric-val" id="calc-priority-val" style="color: var(--accent-cyan);">MEDIUM</span>
                </div>
            </div>
        </div>
    </div>

    <!-- Detail Modal Dialog -->
    <div class="modal-overlay" id="modal-overlay">
        <div class="modal-content">
            <button class="modal-close" id="modal-close">&times;</button>
            <div class="modal-header">
                <h2 id="m-name">Business Name</h2>
                <div id="m-category" class="modal-category">RESTAURANTS & CAFES</div>
            </div>

            <div class="modal-section">
                <div class="modal-section-title">Security Profile Details</div>
                <div class="details-grid">
                    <div class="detail-item">
                        <div class="detail-label">Security Audit Score</div>
                        <div class="detail-value" id="m-score">3 / 5</div>
                    </div>
                    <div class="detail-item">
                        <div class="detail-label">Risk Posture</div>
                        <div class="detail-value" id="m-risk">Moderate Risk</div>
                    </div>
                    <div class="detail-item">
                        <div class="detail-label">Outreach Urgency</div>
                        <div class="detail-value" id="m-priority">CRITICAL</div>
                    </div>
                    <div class="detail-item">
                        <div class="detail-label">Maps Rating & Reviews</div>
                        <div class="detail-value" id="m-reviews">4.5 (1,850 Reviews)</div>
                    </div>
                </div>
            </div>

            <div class="modal-section">
                <div class="modal-section-title">Web Server Protocol Checks</div>
                <div class="details-grid">
                    <div class="detail-item">
                        <div class="detail-label">HTTPS Redirection</div>
                        <div class="detail-value" id="m-https">Yes</div>
                    </div>
                    <div class="detail-item">
                        <div class="detail-label">SSL Certificate</div>
                        <div class="detail-value" id="m-ssl">Valid</div>
                    </div>
                    <div class="detail-item">
                        <div class="detail-label">Admin Exposed</div>
                        <div class="detail-value" id="m-admin">No</div>
                    </div>
                    <div class="detail-item">
                        <div class="detail-label">Safe Browsing Status</div>
                        <div class="detail-value" id="m-safe">Clean</div>
                    </div>
                </div>
            </div>

            <div class="modal-section">
                <div class="modal-section-title">Key Security Vulnerability Flags</div>
                <div class="vulnerability-badge-list" id="m-flags">
                    <!-- Populated dynamically -->
                </div>
            </div>

            <div class="modal-section">
                <div class="modal-section-title">Sales Outreach Pitch Hook</div>
                <div class="pitch-box" id="m-pitch">
                    Sales hook details.
                </div>
            </div>

            <div class="modal-section">
                <div class="modal-section-title">Actionable Remediation Checklist</div>
                <div class="recs-box" id="m-recs">
                    Recommendations details.
                </div>
            </div>
        </div>
    </div>

    <!-- Data & Interactive Logic Scripts -->
    <script>
        // Embed the 30 real business database structures
        const data = {json.dumps(processed_data)};

        // Active State variables
        let selectedCategory = 'all';
        let selectedPriority = 'all';
        let searchQuery = '';
        let viewMode = 'cards'; // 'cards' or 'table'

        // DOM elements
        const leadsGrid = document.getElementById('leads-grid');
        const tableContainer = document.getElementById('table-container');
        const tableBody = document.getElementById('table-body');
        const searchInput = document.getElementById('search-input');
        const emptyState = document.getElementById('empty-state');
        
        // Modal DOM elements
        const modalOverlay = document.getElementById('modal-overlay');
        const modalClose = document.getElementById('modal-close');
        
        // Render Functions
        function renderDashboard() {
            // Apply filtering logic
            const filtered = data.filter(item => {
                const matchesCategory = selectedCategory === 'all' || item.category === selectedCategory;
                const matchesPriority = selectedPriority === 'all' || item.outreach_priority === selectedPriority;
                
                const searchLower = searchQuery.toLowerCase();
                const matchesSearch = searchQuery === '' || 
                    item.name.toLowerCase().includes(searchLower) ||
                    item.location.toLowerCase().includes(searchLower) ||
                    item.key_security_gaps.toLowerCase().includes(searchLower) ||
                    item.outreach_priority.toLowerCase().includes(searchLower) ||
                    item.risk_level.toLowerCase().includes(searchLower);
                
                return matchesCategory && matchesPriority && matchesSearch;
            });

            // Update KPIs based on current filtered view
            updateKPIs(filtered);
            
            // Draw interactive SVG charts
            drawDonutChart(filtered);

            // Toggle empty state warning
            if (filtered.length === 0) {
                emptyState.style.display = 'block';
                leadsGrid.style.display = 'none';
                tableContainer.style.display = 'none';
                return;
            } else {
                emptyState.style.display = 'none';
                if (viewMode === 'cards') {
                    leadsGrid.style.display = 'grid';
                    tableContainer.style.display = 'none';
                } else {
                    leadsGrid.style.display = 'none';
                    tableContainer.style.display = 'block';
                }
            }

            // Render Cards View
            leadsGrid.innerHTML = '';
            filtered.forEach(item => {
                const card = document.createElement('div');
                card.className = 'lead-card';
                card.addEventListener('click', () => openModal(item));

                let priorityClass = item.outreach_priority.toLowerCase();
                let scoreClass = 'secure';
                if (item.security_score <= 2) scoreClass = 'danger';
                else if (item.security_score === 3) scoreClass = 'warning';

                card.innerHTML = `
                    <div class="card-header">
                        <div class="card-title">
                            <h3>${item.name}</h3>
                            <p>${item.location}</p>
                        </div>
                        <div class="score-badge">
                            <span class="score-val ${scoreClass}">${item.security_score}/5</span>
                            <span class="score-label">SCORE</span>
                        </div>
                    </div>
                    <div class="card-body">
                        <div class="meta-item">
                            <span class="meta-label">Category:</span>
                            <span class="meta-val">${item.category}</span>
                        </div>
                        <div class="meta-item">
                            <span class="meta-label">Website:</span>
                            <span class="meta-val">${item.website_url !== 'None' ? `<a href="${item.website_url}" target="_blank" style="color: var(--accent-cyan); text-decoration: none;" onclick="event.stopPropagation()">${item.website_url.replace('https://','').replace('http://','')}</a>` : '<span style="color: var(--text-muted);">None</span>'}</span>
                        </div>
                        <div class="meta-item">
                            <span class="meta-label">Google Rating:</span>
                            <span class="meta-val">★ ${item.rating} (${item.reviews.toLocaleString()} reviews)</span>
                        </div>
                        <div class="meta-item">
                            <span class="meta-label">Outreach Urgency:</span>
                            <span class="priority-indicator ${priorityClass}">${item.outreach_priority}</span>
                        </div>
                    </div>
                    <div class="card-footer" title="${item.sales_pitch_hook}">
                        Hook: ${item.sales_pitch_hook}
                    </div>
                `;
                leadsGrid.appendChild(card);
            });

            // Render Table View
            tableBody.innerHTML = '';
            filtered.forEach((item, idx) => {
                const row = document.createElement('tr');
                if (idx % 2 === 1) row.className = 'zebra';
                row.addEventListener('click', () => openModal(item));

                let priorityClass = item.outreach_priority.toLowerCase();
                let scoreClass = 'secure';
                if (item.security_score <= 2) scoreClass = 'danger';
                else if (item.security_score === 3) scoreClass = 'warning';

                row.innerHTML = `
                    <td><code>${item.id}</code></td>
                    <td><strong>${item.name}</strong></td>
                    <td>${item.website_url !== 'None' ? `<a href="${item.website_url}" target="_blank" style="color: var(--accent-cyan); text-decoration: none;" onclick="event.stopPropagation()">${item.website_url.replace('https://','').replace('http://','')}</a>` : '<span style="color: var(--text-muted);">None</span>'}</td>
                    <td>${item.category}</td>
                    <td>${item.location}</td>
                    <td><strong class="${scoreClass}">${item.security_score} / 5</strong></td>
                    <td>${item.risk_level}</td>
                    <td><span class="priority-indicator ${priorityClass}">${item.outreach_priority}</span></td>
                    <td>${item.https_redirect}</td>
                    <td>${item.ssl_status}</td>
                `;
                tableBody.appendChild(row);
            });
        }

        function updateKPIs(filteredList) {
            document.getElementById('kpi-total').innerText = filteredList.length;
            document.getElementById('donut-center-count').innerText = filteredList.length;
            
            if (filteredList.length === 0) {
                document.getElementById('kpi-score').innerText = '0.00';
                document.getElementById('kpi-secure').innerText = '0.0%';
                document.getElementById('kpi-risky').innerText = '0';
                return;
            }

            const totalScore = filteredList.reduce((sum, item) => sum + item.security_score, 0);
            const avgScore = totalScore / filteredList.length;
            document.getElementById('kpi-score').innerText = avgScore.toFixed(2);

            const secureCount = filteredList.filter(item => item.security_score >= 4).length;
            const secureRate = (secureCount / filteredList.length) * 100;
            document.getElementById('kpi-secure').innerText = secureRate.toFixed(1) + '%';

            const riskyCount = filteredList.filter(item => item.security_score <= 2).length;
            document.getElementById('kpi-risky').innerText = riskyCount;
        }

        // Draw dynamic SVG Donut Chart
        function drawDonutChart(filteredList) {
            const group = document.getElementById('donut-segments-group');
            group.innerHTML = '';
            
            if (filteredList.length === 0) return;

            // Calculate frequencies
            const counts = {
                "Very Secure": filteredList.filter(item => item.risk_level === 'Very Secure').length,
                "Secure": filteredList.filter(item => item.risk_level === 'Secure').length,
                "High Risk": filteredList.filter(item => item.risk_level === 'High Risk').length,
                "Very High Risk": filteredList.filter(item => item.risk_level === 'Very High Risk').length
            };

            const colors = {
                "Very Secure": "var(--c-very-secure)",
                "Secure": "var(--c-secure)",
                "High Risk": "var(--c-high)",
                "Very High Risk": "var(--c-very-high)"
            };

            const total = filteredList.length;
            let currentOffset = 0;
            const r = 60;
            const circumference = 2 * Math.PI * r; // 376.991

            Object.keys(counts).forEach(key => {
                const val = counts[key];
                if (val === 0) return;

                const percent = val / total;
                const strokeLength = percent * circumference;
                const strokeOffset = circumference - currentOffset;

                const segment = document.createElementNS("http://www.w3.org/2000/svg", "circle");
                segment.setAttribute("class", "donut-segment");
                segment.setAttribute("cx", "80");
                segment.setAttribute("cy", "80");
                segment.setAttribute("r", r.toString());
                segment.setAttribute("stroke", colors[key]);
                segment.setAttribute("stroke-dasharray", `${strokeLength} ${circumference}`);
                segment.setAttribute("stroke-dashoffset", strokeOffset.toString());
                
                group.appendChild(segment);
                currentOffset += strokeLength;
            });
        }

        // Modal Logic
        function openModal(item) {
            document.getElementById('m-name').innerHTML = `${item.name} ${item.website_url !== 'None' ? `<a href="${item.website_url}" target="_blank" style="font-size: 13px; margin-left: 12px; color: var(--accent-cyan); text-decoration: none; border: 1px solid rgba(0, 240, 255, 0.3); padding: 4px 10px; border-radius: 6px; background: rgba(0, 240, 255, 0.05); font-weight: 600; display: inline-flex; align-items: center;" onclick="event.stopPropagation()">🔗 Visit Website</a>` : ''}`;
            document.getElementById('m-category').innerText = item.category;
            
            let scoreClass = 'secure';
            if (item.security_score <= 2) scoreClass = 'red';
            else if (item.security_score === 3) scoreClass = 'yellow';
            
            document.getElementById('m-score').className = `detail-value ${scoreClass}`;
            document.getElementById('m-score').innerText = `${item.security_score} / 5`;
            
            document.getElementById('m-risk').innerText = item.risk_level;
            
            let pClass = 'detail-value red';
            if (item.outreach_priority === 'LOW') pClass = 'detail-value green';
            document.getElementById('m-priority').className = pClass;
            document.getElementById('m-priority').innerText = item.outreach_priority;
            
            document.getElementById('m-reviews').innerText = `★ ${item.rating} (${item.reviews.toLocaleString()} Reviews)`;
            
            document.getElementById('m-https').innerText = item.https_redirect;
            document.getElementById('m-https').className = item.https_redirect === 'Yes' ? 'detail-value green' : 'detail-value red';
            
            document.getElementById('m-ssl').innerText = item.ssl_status;
            document.getElementById('m-ssl').className = item.ssl_status === 'Valid' ? 'detail-value green' : 'detail-value red';
            
            document.getElementById('m-admin').innerText = item.admin_panel_exposed;
            document.getElementById('m-admin').className = item.admin_panel_exposed === 'Yes' ? 'detail-value red' : 'detail-value green';
            
            document.getElementById('m-safe').innerText = item.safe_browsing;
            document.getElementById('m-safe').className = item.safe_browsing === 'Clean' ? 'detail-value green' : 'detail-value red';
            
            // Build flags list
            const flagsContainer = document.getElementById('m-flags');
            flagsContainer.innerHTML = '';
            if (item.key_security_gaps === 'None (All Secure)') {
                flagsContainer.innerHTML = `<span class="vulnerability-tag secure">✔ ISO/OWASP Posture Compliant (No Active Scans Blocked)</span>`;
            } else {
                const gapsArray = item.key_security_gaps.split(', ');
                gapsArray.forEach(gap => {
                    flagsContainer.innerHTML += `<span class="vulnerability-tag">⚠ ${gap}</span>`;
                });
            }
            
            document.getElementById('m-pitch').innerText = item.sales_pitch_hook;
            document.getElementById('m-recs').innerText = item.recommendation;
            
            modalOverlay.classList.add('active');
        }

        function closeModal() {
            modalOverlay.classList.remove('active');
        }

        // Event Listeners
        searchInput.addEventListener('input', (e) => {
            searchQuery = e.target.value;
            renderDashboard();
        });

        // Category filter clicks
        const catButtons = document.querySelectorAll('#filter-category .btn');
        catButtons.forEach(btn => {
            btn.addEventListener('click', (e) => {
                catButtons.forEach(b => b.classList.remove('active'));
                btn.classList.add('active');
                selectedCategory = btn.getAttribute('data-cat');
                renderDashboard();
            });
        });

        // Priority filter clicks
        const priButtons = document.querySelectorAll('#filter-priority .btn');
        priButtons.forEach(btn => {
            btn.addEventListener('click', (e) => {
                priButtons.forEach(b => b.classList.remove('active'));
                btn.classList.add('active');
                selectedPriority = btn.getAttribute('data-priority');
                renderDashboard();
            });
        });

        // View Switcher toggles
        document.getElementById('view-cards-btn').addEventListener('click', () => {
            document.getElementById('view-cards-btn').classList.add('active');
            document.getElementById('view-table-btn').classList.remove('active');
            viewMode = 'cards';
            renderDashboard();
        });

        document.getElementById('view-table-btn').addEventListener('click', () => {
            document.getElementById('view-table-btn').classList.add('active');
            document.getElementById('view-cards-btn').classList.remove('active');
            viewMode = 'table';
            renderDashboard();
        });

        modalClose.addEventListener('click', closeModal);
        modalOverlay.addEventListener('click', (e) => {
            if (e.target === modalOverlay) closeModal();
        });

        // Escape key to close modal
        document.addEventListener('keydown', (e) => {
            if (e.key === 'Escape') closeModal();
        });


        // Category scores bar chart calculation & drawing
        function drawCategoryBarChart(filteredList) {
            const container = document.getElementById('bar-chart-container');
            container.innerHTML = '';
            
            if (filteredList.length === 0) return;

            // Categories list
            const categories = [
                "Hospitality & Hotels",
                "Healthcare & Clinics",
                "Professional Services",
                "Retail & Shops",
                "Restaurants & Cafes"
            ];

            categories.forEach(cat => {
                const catItems = filteredList.filter(item => item.category === cat);
                let avg = 0;
                if (catItems.length > 0) {
                    const totalScore = catItems.reduce((sum, item) => sum + item.security_score, 0);
                    avg = totalScore / catItems.length;
                }

                // Width percentage (score is 0-5)
                const widthPct = (avg / 5) * 100;

                const row = document.createElement('div');
                row.className = 'bar-row';
                
                // Short label
                let shortLabel = cat;
                if (cat === "Professional Services") shortLabel = "Prof. Services";
                else if (cat === "Healthcare & Clinics") shortLabel = "Healthcare";
                else if (cat === "Restaurants & Cafes") shortLabel = "Restaurants";
                else if (cat === "Hospitality & Hotels") shortLabel = "Hospitality";
                else if (cat === "Retail & Shops") shortLabel = "Retail";

                row.innerHTML = `
                    <div class="bar-label" title="${cat}">${shortLabel}</div>
                    <div class="bar-track">
                        <div class="bar-fill" style="width: ${widthPct}%"></div>
                    </div>
                    <div class="bar-val">${avg.toFixed(1)}</div>
                `;
                container.appendChild(row);
            });
        }

        // CSV Export Trigger
        document.getElementById('export-csv-btn').addEventListener('click', () => {
            const headers = ["Business ID", "Business Name", "Category", "Location", "Security Score", "Risk Level", "GBP Claimed", "Website Status", "Website URL", "Outreach Priority", "Estimated Lead Value", "Sales Pitch Hook", "Primary Security Gaps", "Actionable Recommendation"];
            
            // Get filtered list matching active controls
            const filtered = data.filter(item => {
                const matchesCategory = selectedCategory === 'all' || item.category === selectedCategory;
                const matchesPriority = selectedPriority === 'all' || item.outreach_priority === selectedPriority;
                
                const searchLower = searchQuery.toLowerCase();
                const matchesSearch = searchQuery === '' || 
                    item.name.toLowerCase().includes(searchLower) ||
                    item.location.toLowerCase().includes(searchLower) ||
                    item.key_security_gaps.toLowerCase().includes(searchLower) ||
                    item.outreach_priority.toLowerCase().includes(searchLower) ||
                    item.risk_level.toLowerCase().includes(searchLower);
                
                return matchesCategory && matchesPriority && matchesSearch;
            });

            let csvContent = "data:text/csv;charset=utf-8,\ufeff";
            csvContent += headers.join(",") + "\n";

            filtered.forEach(item => {
                const row = [
                    item.id,
                    `"${item.name.replace(/"/g, '""')}"`,
                    item.category,
                    item.location,
                    item.security_score,
                    item.risk_level,
                    item.gbp_claimed,
                    item.website_status,
                    item.website_url,
                    item.outreach_priority,
                    `"${item.estimated_lead_value.replace(/"/g, '""')}"`,
                    `"${item.sales_pitch_hook.replace(/"/g, '""')}"`,
                    `"${item.key_security_gaps.replace(/"/g, '""')}"`,
                    `"${item.recommendation.replace(/"/g, '""')}"`
                ];
                csvContent += row.join(",") + "\n";
            });

            const encodedUri = encodeURI(csvContent);
            const link = document.createElement("a");
            link.setAttribute("href", encodedUri);
            link.setAttribute("download", `Filtered_Security_Audit_Leads_${new Date().toISOString().slice(0,10)}.csv`);
            document.body.appendChild(link);
            link.click();
            document.body.removeChild(link);
        });

        // Tab Navigation Event Listeners
        const tabLeadsBtn = document.getElementById('tab-leads-btn');
        const tabReportBtn = document.getElementById('tab-report-btn');
        const leadsSection = document.getElementById('leads-section-container');
        const reportSection = document.getElementById('report-section-container');
        const dashboardTopGrid = document.getElementById('dashboard-top-grid');

        tabLeadsBtn.addEventListener('click', () => {
            tabLeadsBtn.classList.add('active');
            tabReportBtn.classList.remove('active');
            leadsSection.style.display = 'block';
            reportSection.style.display = 'none';
            dashboardTopGrid.style.display = 'grid';
        });

        tabReportBtn.addEventListener('click', () => {
            tabReportBtn.classList.add('active');
            tabLeadsBtn.classList.remove('active');
            leadsSection.style.display = 'none';
            reportSection.style.display = 'block';
            dashboardTopGrid.style.display = 'none'; // Hide dynamic mini charts on report tab for clean reading layout
        });

        // Extend renderDashboard to update category bar chart
        const originalRender = renderDashboard;
        renderDashboard = function() {
            originalRender();
            
            // Get filtered list for bar chart
            const filteredForBar = data.filter(item => {
                const matchesCategory = selectedCategory === 'all' || item.category === selectedCategory;
                const matchesPriority = selectedPriority === 'all' || item.outreach_priority === selectedPriority;
                
                const searchLower = searchQuery.toLowerCase();
                const matchesSearch = searchQuery === '' || 
                    item.name.toLowerCase().includes(searchLower) ||
                    item.location.toLowerCase().includes(searchLower) ||
                    item.key_security_gaps.toLowerCase().includes(searchLower) ||
                    item.outreach_priority.toLowerCase().includes(searchLower) ||
                    item.risk_level.toLowerCase().includes(searchLower);
                
                return matchesCategory && matchesPriority && matchesSearch;
            });
            drawCategoryBarChart(filteredForBar);
        };

        // City Filtering Logic
        let selectedCity = 'all';

        function renderCityPills() {
            const row = document.getElementById('city-pills-row');
            
            // Get unique cities
            const citiesSet = new Set(data.map(item => item.location));
            const citiesList = ['all', ...Array.from(citiesSet).sort()];
            
            row.innerHTML = '';
            citiesList.forEach(city => {
                const count = city === 'all' ? data.length : data.filter(item => item.location === city).length;
                const pill = document.createElement('div');
                pill.className = `city-pill ${selectedCity === city ? 'active' : ''}`;
                pill.innerHTML = city === 'all' ? `🌐 All Cities (${count})` : `📍 ${city} (${count})`;
                
                pill.addEventListener('click', () => {
                    selectedCity = city;
                    renderCityPills();
                    renderDashboard();
                });
                row.appendChild(pill);
            });
        }

        // Score Calculator Logic
        const openCalcBtn = document.getElementById('open-calc-btn');
        const closeCalcBtn = document.getElementById('close-calc-btn');
        const calcDrawer = document.getElementById('calc-drawer');
        const drawerOverlay = document.getElementById('drawer-overlay');

        openCalcBtn.addEventListener('click', () => {
            drawerOverlay.style.display = 'block';
            setTimeout(() => {
                drawerOverlay.style.opacity = '1';
                calcDrawer.classList.add('open');
            }, 10);
        });

        function closeDrawer() {
            calcDrawer.classList.remove('open');
            drawerOverlay.style.opacity = '0';
            setTimeout(() => {
                drawerOverlay.style.display = 'none';
            }, 300);
        }

        closeCalcBtn.addEventListener('click', closeDrawer);
        drawerOverlay.addEventListener('click', closeDrawer);

        // Score Meter Calculation
        const calcHttp = document.getElementById('calc-v-http');
        const calcSsl = document.getElementById('calc-v-ssl');
        const calcAdmin = document.getElementById('calc-v-admin');
        const calcForms = document.getElementById('calc-v-forms');
        const calcSafe = document.getElementById('calc-v-safe');

        function recomputeSimulatedScore() {
            let score = 5.0;
            if (calcHttp.checked) score -= 1.0;
            if (calcSsl.checked) score -= 1.0;
            if (calcAdmin.checked) score -= 1.0;
            if (calcForms.checked) score -= 1.0;
            if (calcSafe.checked) score -= 1.0;
            
            if (score < 1.0) score = 1.0;

            document.getElementById('calc-score-val').textContent = score.toFixed(1);

            const riskVal = document.getElementById('calc-risk-val');
            const priorityVal = document.getElementById('calc-priority-val');

            if (score >= 5.0) {
                riskVal.textContent = "Very Secure";
                riskVal.style.color = "var(--color-secure)";
                priorityVal.textContent = "LOW";
                priorityVal.style.color = "var(--color-secure)";
            } else if (score >= 4.0) {
                riskVal.textContent = "Secure";
                riskVal.style.color = "var(--color-secure)";
                priorityVal.textContent = "LOW";
                priorityVal.style.color = "var(--color-secure)";
            } else if (score >= 3.0) {
                riskVal.textContent = "Moderate Risk";
                riskVal.style.color = "var(--accent-cyan)";
                priorityVal.textContent = "MEDIUM";
                priorityVal.style.color = "var(--accent-cyan)";
            } else if (score >= 2.0) {
                riskVal.textContent = "High Risk";
                riskVal.style.color = "var(--color-danger)";
                priorityVal.textContent = "CRITICAL";
                priorityVal.style.color = "var(--color-danger)";
            } else {
                riskVal.textContent = "Very High Risk";
                riskVal.style.color = "var(--color-danger)";
                priorityVal.textContent = "CRITICAL";
                priorityVal.style.color = "var(--color-danger)";
            }
        }

        [calcHttp, calcSsl, calcAdmin, calcForms, calcSafe].forEach(el => {
            el.addEventListener('change', recomputeSimulatedScore);
        });

        // Initialize score once
        recomputeSimulatedScore();

        // Copy Email Pitch Clipboard Trigger
        let currentEmailPitchText = '';
        document.getElementById('copy-email-btn').addEventListener('click', () => {
            if (!currentEmailPitchText) return;
            navigator.clipboard.writeText(currentEmailPitchText).then(() => {
                const btn = document.getElementById('copy-email-btn');
                const oldText = btn.textContent;
                btn.textContent = 'Copied! ✓';
                btn.style.color = 'var(--color-secure)';
                btn.style.borderColor = 'var(--color-secure)';
                setTimeout(() => {
                    btn.textContent = oldText;
                    btn.style.color = 'var(--accent-cyan)';
                    btn.style.borderColor = 'var(--accent-cyan)';
                }, 2000);
            });
        });

        // Override lead click detailed modal to populate pitch email
        const originalShowDetail = showDetail;
        showDetail = function(id) {
            originalShowDetail(id);
            
            const lead = data.find(item => item.id === id);
            if (lead) {
                // Populate Email Pitch Text
                const emailText = `Subject: ${lead.outreach_subject_line}

Dear Management Team at ${lead.name},

My name is Audit Specialist, and I recently completed a digital vulnerability review of businesses located in ${lead.location}.

While evaluating your online portal (${lead.website_url}), we identified several core security risks:
* Gaps Detected: ${lead.key_security_gaps}

Specifically: ${lead.first_touch_pitch}

We help companies in the ${lead.category} sector harden these configurations to block competitor hijacking and protect customer forms. We can remediate these items for your team today.

Are you available for a brief 10-minute call this Thursday to verify your settings?

Best regards,
Lead Security Auditor
National SMB Digital Risk Auditor`;

                currentEmailPitchText = emailText;
                document.getElementById('modal-email-pitch').textContent = emailText;
            }
        };

        // Override renderDashboard to filter by selectedCity as well!
        renderDashboard = function() {
            // Update KPIs
            const filteredForKPIs = data.filter(item => {
                const matchesCategory = selectedCategory === 'all' || item.category === selectedCategory;
                const matchesPriority = selectedPriority === 'all' || item.outreach_priority === selectedPriority;
                const matchesCity = selectedCity === 'all' || item.location === selectedCity;
                const searchLower = searchQuery.toLowerCase();
                const matchesSearch = searchQuery === '' || 
                    item.name.toLowerCase().includes(searchLower) ||
                    item.location.toLowerCase().includes(searchLower) ||
                    item.key_security_gaps.toLowerCase().includes(searchLower) ||
                    item.outreach_priority.toLowerCase().includes(searchLower) ||
                    item.risk_level.toLowerCase().includes(searchLower);
                
                return matchesCategory && matchesPriority && matchesCity && matchesSearch;
            });

            // Update KPI values
            document.getElementById('kpi-total').textContent = filteredForKPIs.length;
            
            let avgScore = 0;
            if (filteredForKPIs.length > 0) {
                const totalScore = filteredForKPIs.reduce((sum, item) => sum + item.security_score, 0);
                avgScore = totalScore / filteredForKPIs.length;
            }
            document.getElementById('kpi-score').textContent = avgScore.toFixed(2);

            const secureCount = filteredForKPIs.filter(item => item.security_score >= 4).length;
            const secureRate = filteredForKPIs.length > 0 ? (secureCount / filteredForKPIs.length) * 100 : 0;
            document.getElementById('kpi-secure').textContent = secureRate.toFixed(1) + '%';

            const riskyCount = filteredForKPIs.filter(item => item.security_score <= 2).length;
            document.getElementById('kpi-risky').textContent = riskyCount;
            
            // Center Donut Number
            document.getElementById('donut-center-count').textContent = filteredForKPIs.length;

            // Render Ring Slices
            const vSecureCount = filteredForKPIs.filter(item => item.risk_level === 'Very Secure').length;
            const secureLvlCount = filteredForKPIs.filter(item => item.risk_level === 'Secure').length;
            const highCount = filteredForKPIs.filter(item => item.security_score === 2).length;
            const vHighCount = filteredForKPIs.filter(item => item.security_score <= 1).length;

            const totalSegments = vSecureCount + secureLvlCount + highCount + vHighCount;
            const group = document.getElementById('donut-segments-group');
            group.innerHTML = '';

            if (totalSegments > 0) {
                const segments = [
                    { count: vSecureCount, color: 'var(--c-very-secure)' },
                    { count: secureLvlCount, color: 'var(--c-secure)' },
                    { count: highCount, color: 'var(--c-high)' },
                    { count: vHighCount, color: 'var(--c-very-high)' }
                ].filter(s => s.count > 0);

                let accumulatedPct = 0;
                const r = 60;
                const circumference = 2 * Math.PI * r;

                segments.forEach(seg => {
                    const pct = seg.count / totalSegments;
                    const strokeDasharray = `${pct * circumference} ${circumference}`;
                    const strokeDashoffset = -accumulatedPct * circumference;
                    
                    const circle = document.createElementNS('http://www.w3.org/2000/svg', 'circle');
                    circle.setAttribute('class', 'donut-segment');
                    circle.setAttribute('cx', '80');
                    circle.setAttribute('cy', '80');
                    circle.setAttribute('r', r.toString());
                    circle.setAttribute('stroke', seg.color);
                    circle.setAttribute('stroke-dasharray', strokeDasharray);
                    circle.setAttribute('stroke-dashoffset', strokeDashoffset.toString());
                    
                    group.appendChild(circle);
                    accumulatedPct += pct;
                });
            }

            
            const sortVal = document.getElementById('sort-filter') ? document.getElementById('sort-filter').value : 'score-asc';
            if (sortVal === 'score-asc') {
                filteredForKPIs.sort((a, b) => a.security_score - b.security_score || b.reviews - a.reviews);
            } else if (sortVal === 'score-desc') {
                filteredForKPIs.sort((a, b) => b.security_score - a.security_score || b.reviews - a.reviews);
            } else if (sortVal === 'name-asc') {
                filteredForKPIs.sort((a, b) => a.name.localeCompare(b.name));
            } else if (sortVal === 'reviews-desc') {
                filteredForKPIs.sort((a, b) => b.reviews - a.reviews);
            }

            // Draw Category Bars
            drawCategoryBarChart(filteredForKPIs);

            // Render cards or table based on selected view
            if (activeView === 'cards') {
                leadsGrid.style.display = 'grid';
                document.getElementById('table-container').style.display = 'none';
                
                leadsGrid.innerHTML = '';
                if (filteredForKPIs.length === 0) {
                    leadsGrid.innerHTML = '<div style="grid-column: 1/-1; text-align: center; color: var(--text-muted); padding: 40px;">No leads found matching active criteria.</div>';
                    return;
                }

                filteredForKPIs.forEach(lead => {
                    const card = document.createElement('div');
                    card.className = `lead-card priority-${lead.outreach_priority.toLowerCase()}`;
                    card.setAttribute('onclick', `showDetail('${lead.id}')`);
                    
                    const isSecure = lead.security_score >= 4;
                    const ratingStars = '★'.repeat(Math.round(lead.rating)) + '☆'.repeat(5 - Math.round(lead.rating));

                    card.innerHTML = `
                        <div class="card-header">
                            <span class="lead-id">${lead.id}</span>
                            <span class="priority-badge priority-${lead.outreach_priority.toLowerCase()}">${lead.outreach_priority}</span>
                        </div>
                        <h3 class="lead-name">${lead.name}</h3>
                        <div class="lead-meta-row">
                            <span>📍 ${lead.location}</span>
                            <span>🏷 ${lead.category}</span>
                        </div>
                        <div class="rating-row" style="font-size: 11px; margin-bottom: 12px; color: var(--accent-cyan); display: flex; gap: 6px;">
                            <span>${ratingStars}</span>
                            <span style="color: var(--text-muted)">(${lead.reviews} reviews)</span>
                        </div>
                        <div style="display: flex; justify-content: space-between; align-items: center; margin-top: auto; padding-top: 12px; border-top: 1px solid var(--border-muted);">
                            <div style="display: flex; flex-direction: column;">
                                <span style="font-size: 10px; color: var(--text-muted); text-transform: uppercase;">Security Score</span>
                                <span style="font-size: 16px; font-weight: 800; color: ${isSecure ? 'var(--color-secure)' : 'var(--color-danger)'}">${lead.security_score}/5</span>
                            </div>
                            <span class="risk-level-badge ${lead.risk_level.toLowerCase().replace(/ /g, '-')}">${lead.risk_level}</span>
                        </div>
                    `;
                    leadsGrid.appendChild(card);
                });
            } else {
                leadsGrid.style.display = 'none';
                const tableContainer = document.getElementById('table-container');
                tableContainer.style.display = 'block';
                
                const tbody = document.getElementById('leads-tbody');
                tbody.innerHTML = '';

                if (filteredForKPIs.length === 0) {
                    tbody.innerHTML = '<tr><td colspan="7" style="text-align: center; padding: 40px; color: var(--text-muted);">No leads found matching active criteria.</td></tr>';
                    return;
                }

                filteredForKPIs.forEach(lead => {
                    const tr = document.createElement('tr');
                    tr.setAttribute('onclick', `showDetail('${lead.id}')`);
                    
                    const isSecure = lead.security_score >= 4;

                    tr.innerHTML = `
                        <td><strong>${lead.id}</strong></td>
                        <td>
                            <div style="font-weight: 600; color: #ffffff;">${lead.name}</div>
                            <div style="font-size: 10.5px; color: var(--text-muted);">${lead.category}</div>
                        </td>
                        <td>📍 ${lead.location}</td>
                        <td style="text-align: center;">
                            <span style="font-weight: 800; color: ${isSecure ? 'var(--color-secure)' : 'var(--color-danger)'}">${lead.security_score}/5</span>
                        </td>
                        <td style="text-align: center;">
                            <span class="risk-level-badge ${lead.risk_level.toLowerCase().replace(/ /g, '-')}">${lead.risk_level}</span>
                        </td>
                        <td style="text-align: center;">
                            <span class="priority-badge priority-${lead.outreach_priority.toLowerCase()}">${lead.outreach_priority}</span>
                        </td>
                        <td>
                            <div style="font-size: 11.5px; max-width: 250px; overflow: hidden; text-overflow: ellipsis; white-space: nowrap;" title="${lead.key_security_gaps}">
                                ${lead.key_security_gaps}
                            </div>
                        </td>
                    `;
                    tbody.appendChild(tr);
                });
            }
        };

        
        const sortFilter = document.getElementById('sort-filter');
        if(sortFilter) {
            sortFilter.addEventListener('change', renderDashboard);
        }

        
        // Theme Toggle Logic
        const themeBtn = document.getElementById('theme-toggle-btn');
        if (themeBtn) {
            themeBtn.addEventListener('click', () => {
                document.body.classList.toggle('light-mode');
                if (document.body.classList.contains('light-mode')) {
                    themeBtn.textContent = '🌙 Dark Mode';
                    themeBtn.style.color = 'var(--text-light)';
                    themeBtn.style.borderColor = 'var(--text-muted)';
                } else {
                    themeBtn.textContent = '☀️ Light Mode';
                    themeBtn.style.color = 'var(--color-warning)';
                    themeBtn.style.borderColor = 'var(--color-warning)';
                }
            });
        }

        // Initial Load
        renderCityPills();
        renderDashboard();
    </script>
</body>
</html>
"""

with open("National_Digital_Risk_Audit_Dashboard.html", "w", encoding="utf-8") as f:
    f.write(html_template)
print("Dashboard National_Digital_Risk_Audit_Dashboard.html generated successfully!")
